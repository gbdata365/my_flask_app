# -*- coding: utf-8 -*-
"""
기업체현황 분석 조회
====================

주요 기능:
1. 지역별 분석: 권역별/시도별/시군구별 기업체 현황
2. 항목별 분석: 기업규모별/산업분류별/조직형태별/대표자성별/연령그룹별/폐업여부별
   - 항목별 분석 시 total(합계) 컬럼은 제외하고 개별 항목만 표시
3. 조건 필터: 기준년월, 자료유형, 분석유형, 시도 선택
4. DataTables.js 기반 정렬 가능한 테이블 (CDN - 클라우드 배포 호환)
5. Excel/CSV/PPT 다운로드
6. 인사이트 분석 및 요약

필터:
- 기준년월 (base_ym): 연간, 분기, 월간 데이터의 기준시점
- 자료유형 (data_type): 연간/분기/월간
- 분석유형: 지역별(권역별/시도별/시군구별), 항목별(기업규모별/산업분류별 등)
- 시도 선택: 시군구별 분석 시 필요
"""

import sys
import os
import io
import tempfile
from pathlib import Path
from datetime import datetime
import pandas as pd
from flask import request, Response, make_response
from loguru import logger
from urllib.parse import quote

# 상위 디렉토리 module 추가
sys.path.insert(0, str(Path(__file__).parent.parent.parent))
from module.db import get_db_connection
from module.region_config import get_region_sido_mapping

# PPT 유틸리티
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    PPT_AVAILABLE = True
except ImportError as e:
    PPT_AVAILABLE = False
    logger.warning(f"python-pptx 관련 오류: {e}")


def get_filter_options():
    """
    필터 옵션 조회 (giup_summary + dim_admin_area 테이블 기준)

    Returns:
        dict: {
            'base_ym_list': 기준년월 목록,
            'data_type_list': 자료유형 목록,
            'sido_list': 시도 목록 (sido_nm 기반),
            'region_list': 권역 목록
        }
    """
    conn = get_db_connection()

    # 기준년월 목록 조회 (최신 순)
    query_base_ym = """
        SELECT DISTINCT base_ym, base_ym1, data_type
        FROM giup_summary
        ORDER BY base_ym DESC
    """
    df_base_ym = pd.read_sql(query_base_ym, conn)

    # 자료유형 목록
    query_data_type = """
        SELECT data_type FROM (
            SELECT DISTINCT data_type,
                CASE data_type
                    WHEN '연간' THEN 1
                    WHEN '분기' THEN 2
                    WHEN '월간' THEN 3
                    ELSE 4
                END as sort_order
            FROM giup_summary
        ) t
        ORDER BY sort_order
    """
    df_data_type = pd.read_sql(query_data_type, conn)

    # 시도 목록 조회 (dim_admin_area에서 sido_nm 기반으로 조회)
    query_sido = """
        SELECT DISTINCT sido_nm
        FROM dim_admin_area
        WHERE sido_nm IS NOT NULL
        ORDER BY sido_nm
    """
    df_sido = pd.read_sql(query_sido, conn)

    # 권역 목록 조회 (dim_admin_area에서 region_code, region_nm 조회)
    query_region = """
        SELECT DISTINCT region_code, region_nm
        FROM dim_admin_area
        WHERE region_nm IS NOT NULL
        ORDER BY region_code
    """
    df_region = pd.read_sql(query_region, conn)

    conn.close()

    return {
        'base_ym_list': df_base_ym.to_dict('records'),
        'data_type_list': df_data_type['data_type'].tolist(),
        'sido_list': df_sido['sido_nm'].tolist(),  # 시도명 리스트로 반환
        'region_list': df_region.to_dict('records')
    }


def get_analysis_data(base_ym, data_type, analysis_type, sido_nm=None, sigun_detail=False, item_type=None):
    """
    분석 유형에 따른 데이터 조회
    - dim_admin_area 테이블의 지역 정보 사용 (시군구명으로 조인)
    - 권역별은 dim_admin_area의 region_nm 사용
    - 기업규모별/산업분류별 등 항목단위 분석 지원 (total 제외)

    Args:
        base_ym: 기준년월 (YYYYMM)
        data_type: 자료유형 (annual/quarterly/monthly)
        analysis_type: 분석유형 (권역별/시도별/시군구별/기업규모별/산업분류별/연령그룹별/조직형태별/대표자성별별)
        sido_nm: 시도명 (시군구별 분석 시 필요)
        sigun_detail: 하위시군구 상세 표시 여부 (False=구분없이 합계, True=상세)
        item_type: 항목별 분석 시 상세 유형

    Returns:
        pandas.DataFrame: 분석 결과
    """
    conn = get_db_connection()

    # dim_admin_area에서 시도별/시군구별 고유 정보 서브쿼리
    # (시군구명 기준 DISTINCT로 중복 제거, 정렬용 코드 포함)
    dim_sigungu_subquery = """
        (SELECT DISTINCT sido_nm, sigungu_nm, region_nm, region_code, sigungu_code
         FROM dim_admin_area
         WHERE sido_nm IS NOT NULL AND sigungu_nm IS NOT NULL)
    """

    # 시도별 권역 매핑 (sido_nm → region_nm)
    dim_sido_region_subquery = """
        (SELECT DISTINCT sido_nm, region_nm, region_code
         FROM dim_admin_area
         WHERE sido_nm IS NOT NULL AND region_nm IS NOT NULL)
    """

    # item_type이 있으면 지역별+항목별 동시 조회
    if item_type and analysis_type in ['권역별', '시도별', '시군구별']:
        # 지역 그룹핑 컬럼 결정
        if analysis_type == '권역별':
            region_col = 'COALESCE(d.region_nm, \'기타\')'
            region_label = '"권역"'
            group_by_extra = ", COALESCE(d.region_code, '99')"
            order_by = "COALESCE(d.region_code, '99')"
        elif analysis_type == '시도별':
            region_col = 'COALESCE(d.sido_nm, s.sido_nm)'
            region_label = '"시도"'
            group_by_extra = ""
            order_by = "MIN(COALESCE(SUBSTRING(d.sigungu_code, 1, 2), '99'))"
        else:  # 시군구별
            if sido_nm and ('전북' in sido_nm or '전라북도' in sido_nm):
                sido_filter = "AND (d.sido_nm IN ('전라북도', '전북특별자치도') OR (d.sido_nm IS NULL AND s.sido_nm IN ('전라북도', '전북특별자치도')))"
            elif sido_nm:
                sido_filter = f"AND (d.sido_nm = '{sido_nm}' OR (d.sido_nm IS NULL AND s.sido_nm = '{sido_nm}'))"
            else:
                sido_filter = ""

            # 하위시군구 구분 여부에 따라 그룹핑 컬럼 결정
            if sigun_detail:
                # 상세 표시: 시군구별로 개별 표시
                region_col = 'COALESCE(d.sigungu_nm, s.sigun_nm)'
                region_label = '"시군구"'
                group_by_extra = ""
                order_by = "MIN(COALESCE(d.sigungu_code, '99999'))"
            else:
                # 구분없이 (합계): 상위 시군 단위로 그룹화 (하위구 제거)
                region_col = """CASE
                    WHEN COALESCE(d.sigungu_nm, s.sigun_nm) LIKE '%구'
                         AND COALESCE(d.sigungu_nm, s.sigun_nm) NOT LIKE '%시'
                    THEN REGEXP_REPLACE(COALESCE(d.sigungu_nm, s.sigun_nm), ' [가-힣]+구$', '')
                    ELSE COALESCE(d.sigungu_nm, s.sigun_nm)
                END"""
                region_label = '"시군구"'
                group_by_extra = ""
                order_by = "MIN(COALESCE(SUBSTRING(d.sigungu_code, 1, 4), '9999'))"

        # 항목별 컬럼 결정
        if item_type == '기업규모별':
            item_select = """
                SUM(cs.large_other) AS "기타대기업",
                SUM(cs.mid_large) AS "중견기업",
                SUM(cs.mid) AS "중기업",
                SUM(cs.small) AS "소기업",
                SUM(cs.micro) AS "소상공인",
                SUM(cs.excluded) AS "규모판정제외",
                SUM(cs.sangchul) AS "상출기업"
            """
            item_join = "LEFT JOIN giup_detail_corp_size cs ON s.id = cs.summary_id"
        elif item_type == '산업분류별':
            item_select = """
                SUM(i.ind_a) AS "농림어업",
                SUM(i.ind_b) AS "광업",
                SUM(i.ind_c) AS "제조업",
                SUM(i.ind_d) AS "전기가스",
                SUM(i.ind_e) AS "수도하수",
                SUM(i.ind_f) AS "건설업",
                SUM(i.ind_g) AS "도소매업",
                SUM(i.ind_h) AS "운수창고",
                SUM(i.ind_i) AS "숙박음식",
                SUM(i.ind_j) AS "정보통신",
                SUM(i.ind_k) AS "금융보험",
                SUM(i.ind_l) AS "부동산",
                SUM(i.ind_m) AS "전문과학",
                SUM(i.ind_n) AS "사업시설",
                SUM(i.ind_o) AS "공공행정",
                SUM(i.ind_p) AS "교육서비스",
                SUM(i.ind_q) AS "보건복지",
                SUM(i.ind_r) AS "예술스포츠",
                SUM(i.ind_s) AS "협회개인",
                SUM(i.ind_t) AS "가구내고용",
                SUM(i.ind_u) AS "국제기관"
            """
            item_join = "LEFT JOIN giup_detail_industry i ON s.id = i.summary_id"
        elif item_type == '연령그룹별':
            # 10세 단위로 합쳐서 표시 (원본 데이터가 10세 단위이므로)
            item_select = """
                SUM(ag.age_under_19) AS "19세이하",
                SUM(COALESCE(ag.age_20_early, 0) + COALESCE(ag.age_20_late, 0)) AS "20대",
                SUM(COALESCE(ag.age_30_early, 0) + COALESCE(ag.age_30_late, 0)) AS "30대",
                SUM(COALESCE(ag.age_40_early, 0) + COALESCE(ag.age_40_late, 0)) AS "40대",
                SUM(COALESCE(ag.age_50_early, 0) + COALESCE(ag.age_50_late, 0)) AS "50대",
                SUM(COALESCE(ag.age_60_early, 0) + COALESCE(ag.age_60_late, 0)) AS "60대",
                SUM(COALESCE(ag.age_70_early, 0) + COALESCE(ag.age_70_late, 0)) AS "70대",
                SUM(COALESCE(ag.age_80_early, 0) + COALESCE(ag.age_80_over, 0)) AS "80대이상"
            """
            item_join = "LEFT JOIN giup_detail_age_group ag ON s.id = ag.summary_id"
        elif item_type == '조직형태별':
            item_select = """
                SUM(o.indiv_biz) AS "개인사업체",
                SUM(o.corp) AS "회사법인",
                SUM(o.corp_other) AS "회사이외법인",
                SUM(o.non_corp) AS "비법인단체",
                SUM(o.gov_local) AS "국가지방자치단체"
            """
            item_join = "LEFT JOIN giup_detail_org_type o ON s.id = o.summary_id"
        elif item_type == '대표자성별별':
            item_select = """
                SUM(g.male) AS "남성",
                SUM(g.female) AS "여성",
                ROUND(SUM(g.female)::numeric / NULLIF((SUM(g.male) + SUM(g.female))::numeric, 0) * 100, 2) AS "여성비율(%)"
            """
            item_join = "LEFT JOIN giup_detail_gender g ON s.id = g.summary_id"
        elif item_type == '폐업여부별':
            item_select = """
                SUM(st.active) AS "영업중",
                SUM(st.closed) AS "폐업",
                ROUND(SUM(st.closed)::numeric / NULLIF((SUM(st.active) + SUM(st.closed))::numeric, 0) * 100, 2) AS "폐업률(%)"
            """
            item_join = "LEFT JOIN giup_detail_status st ON s.id = st.summary_id"
        else:
            item_select = "SUM(o.total) AS \"총사업체수\""
            item_join = "LEFT JOIN giup_detail_org_type o ON s.id = o.summary_id"

        # 시군구별 필터 처리
        where_filter = ""
        if analysis_type == '시군구별' and sido_nm:
            if '전북' in sido_nm or '전라북도' in sido_nm:
                where_filter = "AND (d.sido_nm IN ('전라북도', '전북특별자치도') OR (d.sido_nm IS NULL AND s.sido_nm IN ('전라북도', '전북특별자치도')))"
            else:
                where_filter = f"AND (d.sido_nm = '{sido_nm}' OR (d.sido_nm IS NULL AND s.sido_nm = '{sido_nm}'))"

        query = f"""
            SELECT
                {region_col} AS {region_label},
                {item_select}
            FROM giup_summary s
            {item_join}
            LEFT JOIN {dim_sigungu_subquery} d
                ON d.sido_nm = s.sido_nm AND d.sigungu_nm = s.sigun_nm
            WHERE s.base_ym = '{base_ym}' AND s.data_type = '{data_type}'
            {where_filter}
            GROUP BY {region_col}{group_by_extra}
            ORDER BY {order_by}
        """

    elif analysis_type == '권역별':
        # 권역별 집계 - dim_admin_area의 region_nm 사용
        # giup_summary.sigun_nm과 dim_admin_area.sigungu_nm으로 조인
        query = f"""
            SELECT
                COALESCE(d.region_nm, '기타') AS "권역",
                COUNT(DISTINCT d.sigungu_nm) AS "시군구수",
                SUM(o.total) AS "총사업체수",
                SUM(o.corp + COALESCE(o.corp_other, 0)) AS "법인수",
                SUM(o.indiv_biz) AS "개인사업자수",
                SUM(g.male) AS "남성대표",
                SUM(g.female) AS "여성대표",
                SUM(st.closed) AS "폐업사업체",
                ROUND(SUM(st.closed)::numeric / NULLIF(SUM(o.total)::numeric, 0) * 100, 2) AS "폐업률(%)"
            FROM giup_summary s
            LEFT JOIN giup_detail_org_type o ON s.id = o.summary_id
            LEFT JOIN giup_detail_gender g ON s.id = g.summary_id
            LEFT JOIN giup_detail_status st ON s.id = st.summary_id
            LEFT JOIN {dim_sigungu_subquery} d
                ON d.sido_nm = s.sido_nm AND d.sigungu_nm = s.sigun_nm
            WHERE s.base_ym = '{base_ym}' AND s.data_type = '{data_type}'
            GROUP BY COALESCE(d.region_nm, '기타'), COALESCE(d.region_code, '99')
            ORDER BY COALESCE(d.region_code, '99')
        """

    elif analysis_type == '시도별':
        # 시도별 집계 - dim_admin_area의 sido_nm 사용
        query = f"""
            SELECT
                COALESCE(d.sido_nm, s.sido_nm) AS "시도",
                COUNT(DISTINCT COALESCE(d.sigungu_nm, s.sigun_nm)) AS "시군구수",
                SUM(o.total) AS "총사업체수",
                SUM(o.corp + COALESCE(o.corp_other, 0)) AS "법인수",
                SUM(o.indiv_biz) AS "개인사업자수",
                SUM(g.male) AS "남성대표",
                SUM(g.female) AS "여성대표",
                SUM(st.closed) AS "폐업사업체",
                ROUND(SUM(st.closed)::numeric / NULLIF(SUM(o.total)::numeric, 0) * 100, 2) AS "폐업률(%)"
            FROM giup_summary s
            LEFT JOIN giup_detail_org_type o ON s.id = o.summary_id
            LEFT JOIN giup_detail_gender g ON s.id = g.summary_id
            LEFT JOIN giup_detail_status st ON s.id = st.summary_id
            LEFT JOIN {dim_sigungu_subquery} d
                ON d.sido_nm = s.sido_nm AND d.sigungu_nm = s.sigun_nm
            WHERE s.base_ym = '{base_ym}' AND s.data_type = '{data_type}'
            GROUP BY COALESCE(d.sido_nm, s.sido_nm)
            ORDER BY MIN(COALESCE(SUBSTRING(d.sigungu_code, 1, 2), '99'))
        """

    elif analysis_type == '시군구별':
        # 시도 필터 (dim_admin_area의 sido_nm 사용)
        # 전북: 전라북도, 전북특별자치도 둘 다 매칭
        if sido_nm and ('전북' in sido_nm or '전라북도' in sido_nm):
            sido_filter = "AND (d.sido_nm IN ('전라북도', '전북특별자치도') OR (d.sido_nm IS NULL AND s.sido_nm IN ('전라북도', '전북특별자치도')))"
        elif sido_nm:
            sido_filter = f"AND (d.sido_nm = '{sido_nm}' OR (d.sido_nm IS NULL AND s.sido_nm = '{sido_nm}'))"
        else:
            sido_filter = ""

        if sigun_detail:
            # 하위시군구 상세 표시 - dim_admin_area의 sigungu_nm 사용
            query = f"""
                SELECT
                    COALESCE(d.sido_nm, s.sido_nm) AS "시도",
                    COALESCE(d.sigungu_nm, s.sigun_nm) AS "시군구명",
                    o.total AS "총사업체수",
                    (o.corp + COALESCE(o.corp_other, 0)) AS "법인수",
                    o.indiv_biz AS "개인사업자수",
                    o.non_corp AS "비법인단체",
                    g.male AS "남성대표",
                    g.female AS "여성대표",
                    st.closed AS "폐업사업체",
                    st.active AS "운영중",
                    ROUND(st.closed::numeric / NULLIF(o.total::numeric, 0) * 100, 2) AS "폐업률(%)"
                FROM giup_summary s
                LEFT JOIN giup_detail_org_type o ON s.id = o.summary_id
                LEFT JOIN giup_detail_gender g ON s.id = g.summary_id
                LEFT JOIN giup_detail_status st ON s.id = st.summary_id
                LEFT JOIN {dim_sigungu_subquery} d
                    ON d.sido_nm = s.sido_nm AND d.sigungu_nm = s.sigun_nm
                WHERE s.base_ym = '{base_ym}' AND s.data_type = '{data_type}'
                {sido_filter}
                ORDER BY COALESCE(d.sigungu_code, '99999')
            """
        else:
            # 하위시군구 구분없이 - 상위 시군 단위로 그룹화
            query = f"""
                SELECT
                    COALESCE(d.sido_nm, s.sido_nm) AS "시도",
                    CASE
                        WHEN COALESCE(d.sigungu_nm, s.sigun_nm) LIKE '%구'
                             AND COALESCE(d.sigungu_nm, s.sigun_nm) NOT LIKE '%시'
                        THEN REGEXP_REPLACE(COALESCE(d.sigungu_nm, s.sigun_nm), ' [가-힣]+구$', '')
                        ELSE COALESCE(d.sigungu_nm, s.sigun_nm)
                    END AS "시군구명",
                    SUM(o.total) AS "총사업체수",
                    SUM(o.corp + COALESCE(o.corp_other, 0)) AS "법인수",
                    SUM(o.indiv_biz) AS "개인사업자수",
                    SUM(o.non_corp) AS "비법인단체",
                    SUM(g.male) AS "남성대표",
                    SUM(g.female) AS "여성대표",
                    SUM(st.closed) AS "폐업사업체",
                    SUM(st.active) AS "운영중",
                    ROUND(SUM(st.closed)::numeric / NULLIF(SUM(o.total)::numeric, 0) * 100, 2) AS "폐업률(%)"
                FROM giup_summary s
                LEFT JOIN giup_detail_org_type o ON s.id = o.summary_id
                LEFT JOIN giup_detail_gender g ON s.id = g.summary_id
                LEFT JOIN giup_detail_status st ON s.id = st.summary_id
                LEFT JOIN {dim_sigungu_subquery} d
                    ON d.sido_nm = s.sido_nm AND d.sigungu_nm = s.sigun_nm
                WHERE s.base_ym = '{base_ym}' AND s.data_type = '{data_type}'
                {sido_filter}
                GROUP BY COALESCE(d.sido_nm, s.sido_nm),
                         CASE
                             WHEN COALESCE(d.sigungu_nm, s.sigun_nm) LIKE '%구'
                                  AND COALESCE(d.sigungu_nm, s.sigun_nm) NOT LIKE '%시'
                             THEN REGEXP_REPLACE(COALESCE(d.sigungu_nm, s.sigun_nm), ' [가-힣]+구$', '')
                             ELSE COALESCE(d.sigungu_nm, s.sigun_nm)
                         END
                ORDER BY MIN(COALESCE(SUBSTRING(d.sigungu_code, 1, 4), '9999'))
            """

    elif analysis_type == '기업규모별':
        # 기업규모별 집계 - total 제외하고 개별 항목만 표시
        query = f"""
            SELECT
                '전국' AS "구분",
                SUM(cs.large_other) AS "기타대기업",
                SUM(cs.mid_large) AS "중견기업",
                SUM(cs.mid) AS "중기업",
                SUM(cs.small) AS "소기업",
                SUM(cs.micro) AS "소상공인",
                SUM(cs.excluded) AS "규모판정제외",
                SUM(cs.sangchul) AS "상출기업"
            FROM giup_summary s
            LEFT JOIN giup_detail_corp_size cs ON s.id = cs.summary_id
            WHERE s.base_ym = '{base_ym}' AND s.data_type = '{data_type}'
        """

    elif analysis_type == '산업분류별':
        # 산업분류별 집계 - total 제외하고 개별 항목만 표시
        query = f"""
            SELECT
                '전국' AS "구분",
                SUM(i.ind_a) AS "농림어업",
                SUM(i.ind_b) AS "광업",
                SUM(i.ind_c) AS "제조업",
                SUM(i.ind_d) AS "전기가스",
                SUM(i.ind_e) AS "수도하수",
                SUM(i.ind_f) AS "건설업",
                SUM(i.ind_g) AS "도소매업",
                SUM(i.ind_h) AS "운수창고",
                SUM(i.ind_i) AS "숙박음식",
                SUM(i.ind_j) AS "정보통신",
                SUM(i.ind_k) AS "금융보험",
                SUM(i.ind_l) AS "부동산",
                SUM(i.ind_m) AS "전문과학",
                SUM(i.ind_n) AS "사업시설",
                SUM(i.ind_o) AS "공공행정",
                SUM(i.ind_p) AS "교육서비스",
                SUM(i.ind_q) AS "보건복지",
                SUM(i.ind_r) AS "예술스포츠",
                SUM(i.ind_s) AS "협회개인",
                SUM(i.ind_t) AS "가구내고용",
                SUM(i.ind_u) AS "국제기관"
            FROM giup_summary s
            LEFT JOIN giup_detail_industry i ON s.id = i.summary_id
            WHERE s.base_ym = '{base_ym}' AND s.data_type = '{data_type}'
        """

    elif analysis_type == '연령그룹별':
        # 연령그룹별 집계 - 10세 단위로 합쳐서 표시 (원본 데이터가 10세 단위이므로)
        query = f"""
            SELECT
                '전국' AS "구분",
                SUM(ag.age_under_19) AS "19세이하",
                SUM(COALESCE(ag.age_20_early, 0) + COALESCE(ag.age_20_late, 0)) AS "20대",
                SUM(COALESCE(ag.age_30_early, 0) + COALESCE(ag.age_30_late, 0)) AS "30대",
                SUM(COALESCE(ag.age_40_early, 0) + COALESCE(ag.age_40_late, 0)) AS "40대",
                SUM(COALESCE(ag.age_50_early, 0) + COALESCE(ag.age_50_late, 0)) AS "50대",
                SUM(COALESCE(ag.age_60_early, 0) + COALESCE(ag.age_60_late, 0)) AS "60대",
                SUM(COALESCE(ag.age_70_early, 0) + COALESCE(ag.age_70_late, 0)) AS "70대",
                SUM(COALESCE(ag.age_80_early, 0) + COALESCE(ag.age_80_over, 0)) AS "80대이상"
            FROM giup_summary s
            LEFT JOIN giup_detail_age_group ag ON s.id = ag.summary_id
            WHERE s.base_ym = '{base_ym}' AND s.data_type = '{data_type}'
        """

    elif analysis_type == '조직형태별':
        # 조직형태별 집계 - total 제외하고 개별 항목만 표시
        query = f"""
            SELECT
                '전국' AS "구분",
                SUM(o.indiv_biz) AS "개인사업체",
                SUM(o.corp) AS "회사법인",
                SUM(o.corp_other) AS "회사이외법인",
                SUM(o.non_corp) AS "비법인단체",
                SUM(o.gov_local) AS "국가지방자치단체"
            FROM giup_summary s
            LEFT JOIN giup_detail_org_type o ON s.id = o.summary_id
            WHERE s.base_ym = '{base_ym}' AND s.data_type = '{data_type}'
        """

    elif analysis_type == '대표자성별별':
        # 대표자성별별 집계 - total 제외하고 개별 항목만 표시
        query = f"""
            SELECT
                '전국' AS "구분",
                SUM(g.male) AS "남성",
                SUM(g.female) AS "여성",
                ROUND(SUM(g.female)::numeric / NULLIF((SUM(g.male) + SUM(g.female))::numeric, 0) * 100, 2) AS "여성비율(%)"
            FROM giup_summary s
            LEFT JOIN giup_detail_gender g ON s.id = g.summary_id
            WHERE s.base_ym = '{base_ym}' AND s.data_type = '{data_type}'
        """

    elif analysis_type == '폐업여부별':
        # 폐업여부별 집계 - total 제외하고 개별 항목만 표시
        query = f"""
            SELECT
                '전국' AS "구분",
                SUM(st.active) AS "영업중",
                SUM(st.closed) AS "폐업",
                ROUND(SUM(st.closed)::numeric / NULLIF((SUM(st.active) + SUM(st.closed))::numeric, 0) * 100, 2) AS "폐업률(%)"
            FROM giup_summary s
            LEFT JOIN giup_detail_status st ON s.id = st.summary_id
            WHERE s.base_ym = '{base_ym}' AND s.data_type = '{data_type}'
        """

    else:
        # 기본: 전체 합계
        query = f"""
            SELECT
                '전국' AS "구분",
                COUNT(DISTINCT COALESCE(d.sigungu_nm, s.sigun_nm)) AS "시군구수",
                SUM(o.total) AS "총사업체수",
                SUM(o.corp + COALESCE(o.corp_other, 0)) AS "법인수",
                SUM(o.indiv_biz) AS "개인사업자수",
                SUM(g.male) AS "남성대표",
                SUM(g.female) AS "여성대표",
                SUM(st.closed) AS "폐업사업체",
                ROUND(SUM(st.closed)::numeric / NULLIF(SUM(o.total)::numeric, 0) * 100, 2) AS "폐업률(%)"
            FROM giup_summary s
            LEFT JOIN giup_detail_org_type o ON s.id = o.summary_id
            LEFT JOIN giup_detail_gender g ON s.id = g.summary_id
            LEFT JOIN giup_detail_status st ON s.id = st.summary_id
            LEFT JOIN {dim_sigungu_subquery} d
                ON d.sido_nm = s.sido_nm AND d.sigungu_nm = s.sigun_nm
            WHERE s.base_ym = '{base_ym}' AND s.data_type = '{data_type}'
        """

    df = pd.read_sql(query, conn)
    conn.close()

    return df


def handle_export(df, export_format, filename_prefix):
    """
    데이터 내보내기 처리

    Args:
        df: 데이터프레임
        export_format: 'excel' 또는 'csv'
        filename_prefix: 파일명 접두어

    Returns:
        Flask Response
    """
    if export_format == 'excel':
        output = io.BytesIO()
        with pd.ExcelWriter(output, engine='openpyxl') as writer:
            df.to_excel(writer, index=False, sheet_name='분석결과')
        output.seek(0)

        filename = f"{filename_prefix}.xlsx"
        encoded_filename = quote(filename)

        response = make_response(output.getvalue())
        response.headers['Content-Type'] = 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
        response.headers['Content-Disposition'] = f"attachment; filename*=UTF-8''{encoded_filename}"
        return response

    elif export_format == 'csv':
        output = io.StringIO()
        df.to_csv(output, index=False, encoding='utf-8-sig')  # BOM for Excel

        filename = f"{filename_prefix}.csv"
        encoded_filename = quote(filename)

        response = make_response(output.getvalue())
        response.headers['Content-Type'] = 'text/csv; charset=utf-8-sig'
        response.headers['Content-Disposition'] = f"attachment; filename*=UTF-8''{encoded_filename}"
        return response

    return None


def generate_insights(df_result, analysis_type, base_ym, data_type, item_type=None):
    """
    데이터 분석 결과에서 인사이트 생성

    Args:
        df_result: 분석 결과 데이터프레임
        analysis_type: 분석유형 (지역구분)
        base_ym: 기준년월
        data_type: 자료유형
        item_type: 항목구분 (선택사항)

    Returns:
        dict: 인사이트 정보
    """
    if df_result is None or len(df_result) == 0:
        return None

    insights = {
        'summary': [],      # 핵심 요약
        'highlights': [],   # 주요 발견사항
        'warnings': [],     # 주의사항
        'statistics': {}    # 통계 수치
    }

    # 지역명 컬럼 확인
    region_col = None
    for col in ['권역', '시도', '시군구', '시군구명', '구분']:
        if col in df_result.columns:
            region_col = col
            break

    if region_col is None:
        return insights

    # 항목별 분석인지 확인 (총사업체수 컬럼이 없으면 항목별)
    is_item_analysis = '총사업체수' not in df_result.columns and item_type

    if is_item_analysis:
        # 항목별 분석: 항목 컬럼들의 합계 계산
        item_cols = [c for c in df_result.columns if c != region_col]
        total_all_items = 0
        numeric_cols = []
        for col in item_cols:
            # 숫자형 컬럼 여부 확인 (다양한 dtype 지원)
            try:
                if pd.api.types.is_numeric_dtype(df_result[col]):
                    numeric_cols.append(col)
                    total_all_items += pd.to_numeric(df_result[col], errors='coerce').fillna(0).sum()
            except Exception:
                pass

        insights['statistics'] = {
            'total_biz': int(total_all_items) if pd.notna(total_all_items) else 0,
            'total_corp': 0,
            'total_indiv': 0,
            'total_closed': 0,
            'corp_ratio': 0,
            'closure_ratio': 0,
            'is_item_analysis': True,  # 항목별 분석 표시
            'item_type': item_type
        }

        # 핵심 요약 - 항목별
        insights['summary'].append(
            f"{data_type} 기준 ({base_ym}) {item_type} 분석 결과입니다."
        )
        insights['summary'].append(
            f"총 {len(df_result)}개 {analysis_type}에서 합계 {int(total_all_items):,}개 사업체가 조회되었습니다."
        )

        # 주요 발견사항 - 항목별 최대값 컬럼 찾기
        if len(numeric_cols) > 0 and len(df_result) > 0:
            col_totals = {}
            for col in numeric_cols:
                try:
                    col_sum = pd.to_numeric(df_result[col], errors='coerce').fillna(0).sum()
                    col_totals[col] = col_sum
                except Exception:
                    pass

            if col_totals:
                max_col = max(col_totals, key=col_totals.get)
                max_val = col_totals[max_col]
                insights['highlights'].append({
                    'icon': '1',
                    'title': f'최다 {item_type.replace("별", "")} 항목',
                    'content': f"{max_col}: {int(max_val):,}개 ({round(max_val/total_all_items*100, 1) if total_all_items > 0 else 0}%)"
                })

                # 지역별 합계 최대 (숫자형 컬럼만 사용)
                numeric_df = df_result[numeric_cols].apply(pd.to_numeric, errors='coerce').fillna(0)
                row_totals = numeric_df.sum(axis=1)
                max_region_idx = row_totals.idxmax()
                max_region_row = df_result.loc[max_region_idx]
                insights['highlights'].append({
                    'icon': '2',
                    'title': f'사업체 최다 {analysis_type.replace("별", "")}',
                    'content': f"{max_region_row[region_col]}: {int(row_totals[max_region_idx]):,}개"
                })
    else:
        # 기본 통계 (지역별 기본 현황)
        total_biz = df_result['총사업체수'].sum() if '총사업체수' in df_result.columns else 0
        total_corp = df_result['법인수'].sum() if '법인수' in df_result.columns else 0
        total_indiv = df_result['개인사업자수'].sum() if '개인사업자수' in df_result.columns else 0
        total_closed = df_result['폐업사업체'].sum() if '폐업사업체' in df_result.columns else 0

        insights['statistics'] = {
            'total_biz': int(total_biz) if pd.notna(total_biz) else 0,
            'total_corp': int(total_corp) if pd.notna(total_corp) else 0,
            'total_indiv': int(total_indiv) if pd.notna(total_indiv) else 0,
            'total_closed': int(total_closed) if pd.notna(total_closed) else 0,
            'corp_ratio': round(total_corp / total_biz * 100, 1) if total_biz > 0 else 0,
            'closure_ratio': round(total_closed / total_biz * 100, 1) if total_biz > 0 else 0
        }

        # 핵심 요약
        insights['summary'].append(
            f"{data_type} 기준 ({base_ym}) 전체 사업체 수는 총 {int(total_biz):,}개입니다."
        )

        if total_corp > 0 and total_indiv > 0:
            corp_pct = round(total_corp / total_biz * 100, 1)
            indiv_pct = round(total_indiv / total_biz * 100, 1)
            insights['summary'].append(
                f"법인 {int(total_corp):,}개({corp_pct}%), 개인사업자 {int(total_indiv):,}개({indiv_pct}%)로 구성됩니다."
            )

    # 주요 발견사항
    if '총사업체수' in df_result.columns:
        # 최대/최소 지역
        max_region = df_result.loc[df_result['총사업체수'].idxmax()]
        min_region = df_result.loc[df_result['총사업체수'].idxmin()]

        insights['highlights'].append({
            'icon': '1',
            'title': '사업체 최다 지역',
            'content': f"{max_region[region_col]}: {int(max_region['총사업체수']):,}개 (전체의 {round(max_region['총사업체수']/total_biz*100, 1)}%)"
        })

        insights['highlights'].append({
            'icon': '2',
            'title': '사업체 최소 지역',
            'content': f"{min_region[region_col]}: {int(min_region['총사업체수']):,}개 (전체의 {round(min_region['총사업체수']/total_biz*100, 1)}%)"
        })

    # 폐업률 분석
    if '폐업률(%)' in df_result.columns:
        avg_closure = df_result['폐업률(%)'].mean()
        max_closure_row = df_result.loc[df_result['폐업률(%)'].idxmax()]
        min_closure_row = df_result.loc[df_result['폐업률(%)'].idxmin()]

        insights['highlights'].append({
            'icon': '3',
            'title': '평균 폐업률',
            'content': f"전체 평균 {avg_closure:.1f}% (최고: {max_closure_row[region_col]} {max_closure_row['폐업률(%)']:.1f}%, 최저: {min_closure_row[region_col]} {min_closure_row['폐업률(%)']:.1f}%)"
        })

        # 폐업률 높은 지역 경고
        high_closure = df_result[df_result['폐업률(%)'] > 10]
        if len(high_closure) > 0:
            regions = ', '.join(high_closure[region_col].tolist()[:3])
            insights['warnings'].append(
                f"폐업률 10% 초과 지역: {regions}" + (f" 외 {len(high_closure)-3}개" if len(high_closure) > 3 else "")
            )

    # 성별 분석
    if '남성대표' in df_result.columns and '여성대표' in df_result.columns:
        total_male = df_result['남성대표'].sum()
        total_female = df_result['여성대표'].sum()
        total_gender = total_male + total_female

        if total_gender > 0:
            female_ratio = round(total_female / total_gender * 100, 1)
            insights['highlights'].append({
                'icon': '4',
                'title': '대표자 성별 비율',
                'content': f"남성 {int(total_male):,}명({round(100-female_ratio, 1)}%), 여성 {int(total_female):,}명({female_ratio}%)"
            })

    return insights


def generate_markdown(df_result, insights, analysis_type, base_ym, data_type, sido_nm=None):
    """
    분석 결과를 마크다운 형식으로 생성

    Args:
        df_result: 분석 결과 데이터프레임
        insights: 인사이트 정보
        analysis_type: 분석유형
        base_ym: 기준년월
        data_type: 자료유형
        sido_nm: 시도명 (시군구별 분석 시)

    Returns:
        str: 마크다운 문서
    """
    now = datetime.now().strftime("%Y-%m-%d %H:%M")

    md = f"""# 기업체현황 분석 보고서

**생성일시:** {now}
**분석유형:** {analysis_type}
**자료유형:** {data_type}
**기준년월:** {base_ym}
{f"**시도:** {sido_nm}" if sido_nm else ""}

---

## 1. 핵심 요약

"""

    if insights and insights.get('summary'):
        for item in insights['summary']:
            md += f"- {item}\n"

    md += "\n## 2. 주요 통계\n\n"

    if insights and insights.get('statistics'):
        stats = insights['statistics']
        md += f"""| 구분 | 수치 |
|------|------|
| 총 사업체수 | {stats.get('total_biz', 0):,}개 |
| 법인 수 | {stats.get('total_corp', 0):,}개 ({stats.get('corp_ratio', 0)}%) |
| 개인사업자 수 | {stats.get('total_indiv', 0):,}개 |
| 폐업 사업체 | {stats.get('total_closed', 0):,}개 ({stats.get('closure_ratio', 0)}%) |

"""

    md += "## 3. 주요 발견사항\n\n"

    if insights and insights.get('highlights'):
        for h in insights['highlights']:
            md += f"### {h['icon']}. {h['title']}\n"
            md += f"{h['content']}\n\n"

    if insights and insights.get('warnings'):
        md += "## 4. 주의사항\n\n"
        for w in insights['warnings']:
            md += f"- {w}\n"

    md += "\n## 5. 상세 데이터\n\n"

    if df_result is not None and len(df_result) > 0:
        # 테이블 헤더
        md += "| " + " | ".join(df_result.columns) + " |\n"
        md += "| " + " | ".join(["---"] * len(df_result.columns)) + " |\n"

        # 테이블 내용 (최대 20행)
        for idx, row in df_result.head(20).iterrows():
            values = []
            for val in row:
                if pd.isna(val):
                    values.append("-")
                elif isinstance(val, float):
                    if val == int(val):
                        values.append(f"{int(val):,}")
                    else:
                        values.append(f"{val:,.2f}")
                elif isinstance(val, int):
                    values.append(f"{val:,}")
                else:
                    values.append(str(val))
            md += "| " + " | ".join(values) + " |\n"

        if len(df_result) > 20:
            md += f"\n*... 외 {len(df_result) - 20}건 생략*\n"

    md += f"""
---

*본 보고서는 기업체현황 분석 시스템에서 자동 생성되었습니다.*
"""

    return md


def generate_ppt(df_result, insights, analysis_type, base_ym, data_type, sido_nm=None):
    """
    PPT 파일 생성

    Args:
        df_result: 분석 결과 데이터프레임
        insights: 인사이트 정보
        analysis_type: 분석유형
        base_ym: 기준년월
        data_type: 자료유형
        sido_nm: 시도명

    Returns:
        str: 생성된 PPT 파일 경로
    """
    if not PPT_AVAILABLE:
        return None

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 제목 슬라이드
    slide_layout = prs.slide_layouts[6]  # 빈 슬라이드
    slide = prs.slides.add_slide(slide_layout)

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = "기업체현황 분석 보고서"
    title_p.font.size = Pt(44)
    title_p.font.bold = True
    title_p.alignment = PP_ALIGN.CENTER

    # 부제목
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    sub_p = subtitle_frame.paragraphs[0]
    region_text = f" - {sido_nm}" if sido_nm else ""
    sub_p.text = f"{analysis_type} 분석 | {data_type} | {base_ym}{region_text}"
    sub_p.font.size = Pt(24)
    sub_p.alignment = PP_ALIGN.CENTER

    # 요약 슬라이드
    slide = prs.slides.add_slide(slide_layout)

    header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    header_frame = header_box.text_frame
    header_p = header_frame.paragraphs[0]
    header_p.text = "핵심 요약"
    header_p.font.size = Pt(32)
    header_p.font.bold = True

    # 통계 카드들
    if insights and insights.get('statistics'):
        stats = insights['statistics']

        cards = [
            ('총 사업체수', f"{stats.get('total_biz', 0):,}개"),
            ('법인', f"{stats.get('total_corp', 0):,}개 ({stats.get('corp_ratio', 0)}%)"),
            ('개인사업자', f"{stats.get('total_indiv', 0):,}개"),
            ('폐업률', f"{stats.get('closure_ratio', 0)}%")
        ]

        for i, (label, value) in enumerate(cards):
            x = Inches(0.5 + (i % 4) * 3.2)
            y = Inches(1.5)

            box = slide.shapes.add_textbox(x, y, Inches(3), Inches(1.5))
            tf = box.text_frame

            p1 = tf.paragraphs[0]
            p1.text = label
            p1.font.size = Pt(16)
            p1.font.bold = True

            p2 = tf.add_paragraph()
            p2.text = value
            p2.font.size = Pt(28)
            p2.font.bold = True

    # 요약 텍스트
    if insights and insights.get('summary'):
        summary_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12.333), Inches(3))
        summary_frame = summary_box.text_frame
        summary_frame.word_wrap = True

        for i, text in enumerate(insights['summary']):
            if i == 0:
                p = summary_frame.paragraphs[0]
            else:
                p = summary_frame.add_paragraph()
            p.text = f"• {text}"
            p.font.size = Pt(18)
            p.space_after = Pt(12)

    # 주요 발견사항 슬라이드
    if insights and insights.get('highlights'):
        slide = prs.slides.add_slide(slide_layout)

        header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
        header_frame = header_box.text_frame
        header_p = header_frame.paragraphs[0]
        header_p.text = "주요 발견사항"
        header_p.font.size = Pt(32)
        header_p.font.bold = True

        y_pos = 1.3
        for h in insights['highlights']:
            box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(12.333), Inches(1.2))
            tf = box.text_frame
            tf.word_wrap = True

            p1 = tf.paragraphs[0]
            p1.text = f"{h['icon']}. {h['title']}"
            p1.font.size = Pt(20)
            p1.font.bold = True

            p2 = tf.add_paragraph()
            p2.text = h['content']
            p2.font.size = Pt(16)

            y_pos += 1.4

    # 데이터 테이블 슬라이드 (상위 10개)
    if df_result is not None and len(df_result) > 0:
        slide = prs.slides.add_slide(slide_layout)

        header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
        header_frame = header_box.text_frame
        header_p = header_frame.paragraphs[0]
        header_p.text = f"상세 데이터 (상위 {min(10, len(df_result))}개)"
        header_p.font.size = Pt(32)
        header_p.font.bold = True

        # 테이블 (최대 10행, 주요 컬럼만)
        display_cols = [c for c in df_result.columns if c not in ['시도코드', '시군구코드']][:6]
        rows = min(10, len(df_result))
        cols = len(display_cols)

        table = slide.shapes.add_table(rows + 1, cols, Inches(0.5), Inches(1.2), Inches(12.333), Inches(5)).table

        # 헤더
        for i, col in enumerate(display_cols):
            cell = table.cell(0, i)
            cell.text = col
            cell.text_frame.paragraphs[0].font.size = Pt(11)
            cell.text_frame.paragraphs[0].font.bold = True

        # 데이터
        for r in range(rows):
            for c, col in enumerate(display_cols):
                cell = table.cell(r + 1, c)
                val = df_result.iloc[r][col]
                if pd.isna(val):
                    cell.text = "-"
                elif isinstance(val, float):
                    cell.text = f"{val:,.1f}" if val != int(val) else f"{int(val):,}"
                elif isinstance(val, int):
                    cell.text = f"{val:,}"
                else:
                    cell.text = str(val)
                cell.text_frame.paragraphs[0].font.size = Pt(10)

    # 파일 저장
    output_path = os.path.join(tempfile.gettempdir(), f"기업체현황_{analysis_type}_{base_ym}.pptx")
    prs.save(output_path)

    return output_path


def handle_ppt_save(form_data, df_result, insights, filter_opts):
    """
    PPT 저장 요청 처리 (MD 파일도 함께 생성)

    Args:
        form_data: 폼 데이터
        df_result: 분석 결과
        insights: 인사이트
        filter_opts: 필터 옵션

    Returns:
        Flask Response
    """
    if not PPT_AVAILABLE:
        return Response(
            "<script>alert('PPT 저장 기능을 사용하려면 python-pptx를 설치하세요.'); history.back();</script>",
            mimetype='text/html'
        )

    try:
        base_ym = form_data.get('base_ym', '')
        data_type = form_data.get('data_type', '')
        analysis_type = form_data.get('analysis_type', '시도별')
        sido_nm = form_data.get('sido_nm', '')  # 시도명 직접 사용

        # MD 파일 생성 (내부 저장)
        md_content = generate_markdown(df_result, insights, analysis_type, base_ym, data_type, sido_nm)
        md_path = os.path.join(tempfile.gettempdir(), f"기업체현황_{analysis_type}_{base_ym}.md")
        with open(md_path, 'w', encoding='utf-8') as f:
            f.write(md_content)
        logger.info(f"MD 파일 생성: {md_path}")

        # PPT 파일 생성
        ppt_path = generate_ppt(df_result, insights, analysis_type, base_ym, data_type, sido_nm)

        if not ppt_path or not os.path.exists(ppt_path):
            return Response(
                "<script>alert('PPT 파일 생성 실패'); history.back();</script>",
                mimetype='text/html'
            )

        # 파일명 생성
        region_suffix = f"_{sido_nm}" if sido_nm else ""
        filename = f"기업체현황_{analysis_type}_{base_ym}{region_suffix}.pptx"

        with open(ppt_path, 'rb') as f:
            file_data = f.read()

        response = make_response(file_data)
        response.headers['Content-Type'] = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
        encoded_filename = quote(filename)
        response.headers['Content-Disposition'] = f"attachment; filename*=UTF-8''{encoded_filename}"
        response.headers['Content-Length'] = len(file_data)

        return response

    except Exception as e:
        import traceback
        logger.error(f"PPT 저장 오류: {str(e)}")
        logger.error(traceback.format_exc())
        return Response(
            f"<script>alert('PPT 저장 중 오류: {str(e)}'); history.back();</script>",
            mimetype='text/html'
        )


def generate_charts(df_result, analysis_type):
    """
    Plotly 차트 생성

    Args:
        df_result: 분석 결과 데이터프레임
        analysis_type: 분석유형

    Returns:
        tuple: (bar_chart_json, pie_chart_json, closure_chart_json)
    """
    import json
    import numpy as np

    # numpy/pandas 타입을 Python 기본 타입으로 변환하는 헬퍼 함수
    def to_python_type(val):
        if isinstance(val, (np.integer, np.int64)):
            return int(val)
        elif isinstance(val, (np.floating, np.float64)):
            return float(val)
        elif pd.isna(val):
            return 0
        return val

    def to_python_list(series):
        return [to_python_type(v) for v in series.fillna(0)]

    if df_result is None or len(df_result) == 0:
        return None, None, None

    # 지역명 컬럼 확인
    region_col = None
    for col in ['권역', '시도', '시군구명', '구분']:
        if col in df_result.columns:
            region_col = col
            break

    if region_col is None:
        return None, None, None

    # 1. 막대차트 - 항목별 분석과 지역별 분석 구분
    if '총사업체수' in df_result.columns:
        # 지역별 분석: 총사업체수 막대차트
        bar_data = {
            'type': 'bar',
            'x': df_result[region_col].tolist(),
            'y': to_python_list(df_result['총사업체수']),
            'marker': {
                'color': 'rgba(52, 152, 219, 0.8)',
                'line': {'color': 'rgba(52, 152, 219, 1)', 'width': 1}
            },
            'text': [f'{int(v):,}' for v in df_result['총사업체수'].fillna(0)],
            'textposition': 'outside',
            'hovertemplate': '%{x}<br>사업체수: %{y:,}<extra></extra>'
        }

        bar_layout = {
            'title': {'text': f'{analysis_type} 사업체 현황', 'font': {'size': 16}},
            'xaxis': {'tickangle': -45, 'tickfont': {'size': 11}},
            'yaxis': {'title': '사업체수', 'tickformat': ','},
            'margin': {'l': 60, 'r': 30, 't': 50, 'b': 100},
            'height': 400,
            'showlegend': False
        }

        bar_chart = {'data': [bar_data], 'layout': bar_layout}
    else:
        # 항목별 분석 (기업규모별, 산업분류별 등): 항목별 막대차트
        # 구분 컬럼 제외한 나머지 컬럼이 항목들
        item_cols = [col for col in df_result.columns if col != region_col]
        if len(item_cols) > 0:
            # 첫 번째 행의 데이터 사용 (전국 합계)
            first_row = df_result.iloc[0]
            x_labels = item_cols
            y_values = [int(first_row[col]) if pd.notna(first_row[col]) else 0 for col in item_cols]

            bar_data = {
                'type': 'bar',
                'x': x_labels,
                'y': y_values,
                'marker': {
                    'color': 'rgba(52, 152, 219, 0.8)',
                    'line': {'color': 'rgba(52, 152, 219, 1)', 'width': 1}
                },
                'text': [f'{v:,}' for v in y_values],
                'textposition': 'outside',
                'hovertemplate': '%{x}<br>사업체수: %{y:,}<extra></extra>'
            }

            bar_layout = {
                'title': {'text': f'{analysis_type} 사업체 현황', 'font': {'size': 16}},
                'xaxis': {'tickangle': -45, 'tickfont': {'size': 11}},
                'yaxis': {'title': '사업체수', 'tickformat': ','},
                'margin': {'l': 60, 'r': 30, 't': 50, 'b': 120},
                'height': 400,
                'showlegend': False
            }

            bar_chart = {'data': [bar_data], 'layout': bar_layout}
        else:
            bar_chart = None

    # 2. 법인/개인 파이차트
    if '법인수' in df_result.columns and '개인사업자수' in df_result.columns:
        corp_total = int(df_result['법인수'].fillna(0).sum())
        indiv_total = int(df_result['개인사업자수'].fillna(0).sum())

        # 법인단체가 있으면 포함
        corp_body_total = 0
        if '법인단체' in df_result.columns:
            corp_body_total = int(df_result['법인단체'].fillna(0).sum())

        pie_labels = ['법인', '개인사업자']
        pie_values = [corp_total, indiv_total]
        pie_colors = ['#3498db', '#e74c3c']

        if corp_body_total > 0:
            pie_labels.append('법인단체')
            pie_values.append(corp_body_total)
            pie_colors.append('#9b59b6')

        pie_data = {
            'type': 'pie',
            'labels': pie_labels,
            'values': pie_values,
            'marker': {'colors': pie_colors},
            'textinfo': 'label+percent',
            'textposition': 'outside',
            'hole': 0.3,
            'hovertemplate': '%{label}<br>%{value:,}개<br>%{percent}<extra></extra>'
        }

        pie_layout = {
            'title': {'text': '조직형태별 구성', 'font': {'size': 16}},
            'height': 400,
            'margin': {'l': 30, 'r': 30, 't': 50, 'b': 30},
            'showlegend': True,
            'legend': {'orientation': 'h', 'y': -0.1}
        }

        pie_chart = {'data': [pie_data], 'layout': pie_layout}
    else:
        pie_chart = None

    # 3. 폐업률 막대차트
    if '폐업률(%)' in df_result.columns:
        closure_y = to_python_list(df_result['폐업률(%)'])
        closure_data = {
            'type': 'bar',
            'x': df_result[region_col].tolist(),
            'y': closure_y,
            'marker': {
                'color': ['rgba(231, 76, 60, 0.8)' if v > 10 else 'rgba(46, 204, 113, 0.8)' for v in closure_y],
                'line': {'width': 1}
            },
            'text': [f'{v:.1f}%' for v in closure_y],
            'textposition': 'outside',
            'hovertemplate': '%{x}<br>폐업률: %{y:.1f}%<extra></extra>'
        }

        closure_layout = {
            'title': {'text': f'{analysis_type} 폐업률 비교', 'font': {'size': 16}},
            'xaxis': {'tickangle': -45, 'tickfont': {'size': 11}},
            'yaxis': {'title': '폐업률 (%)', 'ticksuffix': '%'},
            'margin': {'l': 60, 'r': 30, 't': 50, 'b': 100},
            'height': 400,
            'showlegend': False
        }

        closure_chart = {'data': [closure_data], 'layout': closure_layout}
    else:
        closure_chart = None

    return (
        json.dumps(bar_chart, ensure_ascii=False) if bar_chart else None,
        json.dumps(pie_chart, ensure_ascii=False) if pie_chart else None,
        json.dumps(closure_chart, ensure_ascii=False) if closure_chart else None
    )


def generate_html(filter_opts, selected_base_ym, selected_data_type,
                  selected_region_type, selected_item_type, selected_sido_nm, selected_sigun_detail, df_result, insights=None):
    """
    HTML 생성

    Args:
        filter_opts: 필터 옵션
        selected_base_ym: 선택된 기준년월
        selected_data_type: 선택된 자료유형
        selected_region_type: 선택된 지역구분 (권역별/시도별/시군구별)
        selected_item_type: 선택된 항목구분 (기업규모별/산업분류별 등, 비어있으면 기본현황)
        selected_sido_nm: 선택된 시도명
        selected_sigun_detail: 하위시군구 상세 표시 여부 (False=구분없이, True=상세)
        df_result: 분석 결과 데이터프레임
        insights: 인사이트 정보
    """
    # 분석유형 표시용 문자열
    selected_analysis_type = selected_region_type
    if selected_item_type:
        selected_analysis_type = f"{selected_region_type} ({selected_item_type})"

    # 차트 생성
    bar_chart_json, pie_chart_json, closure_chart_json = generate_charts(df_result, selected_region_type)

    # 인사이트 HTML 생성
    insights_html = ""
    if insights:
        # 핵심 지표 카드
        stats = insights.get('statistics', {})
        is_item_analysis = stats.get('is_item_analysis', False)

        if is_item_analysis:
            # 항목별 분석: 총 사업체수만 표시
            item_type_label = stats.get('item_type', '항목')
            insights_html = f"""
        <div class="insights-section">
            <h2>분석 인사이트 ({item_type_label})</h2>

            <div class="stats-cards">
                <div class="stat-card" style="flex: 2;">
                    <div class="stat-label">총 사업체수 합계</div>
                    <div class="stat-value">{stats.get('total_biz', 0):,}개</div>
                </div>
            </div>

            <div class="summary-box">
                <h3>핵심 요약</h3>
                <ul>
        """
            for s in insights.get('summary', []):
                insights_html += f"<li>{s}</li>"

            insights_html += """
                </ul>
            </div>

            <div class="highlights-grid">
        """

            for h in insights.get('highlights', []):
                insights_html += f"""
                <div class="highlight-card">
                    <div class="highlight-num">{h['icon']}</div>
                    <div class="highlight-content">
                        <div class="highlight-title">{h['title']}</div>
                        <div class="highlight-text">{h['content']}</div>
                    </div>
                </div>
            """

            insights_html += "</div>"  # highlights-grid 닫기

            # 경고사항 (항목별)
            if insights.get('warnings'):
                insights_html += '<div class="warnings-box"><h3>주의사항</h3><ul>'
                for w in insights['warnings']:
                    insights_html += f"<li>{w}</li>"
                insights_html += "</ul></div>"

            insights_html += "</div>"  # insights-section 닫기

        else:
            # 지역별 기본 분석: 전체 카드 표시
            insights_html = f"""
        <div class="insights-section">
            <h2>분석 인사이트</h2>

            <div class="stats-cards">
                <div class="stat-card">
                    <div class="stat-label">총 사업체수</div>
                    <div class="stat-value">{stats.get('total_biz', 0):,}개</div>
                </div>
                <div class="stat-card">
                    <div class="stat-label">법인</div>
                    <div class="stat-value">{stats.get('total_corp', 0):,}개</div>
                    <div class="stat-sub">{stats.get('corp_ratio', 0)}%</div>
                </div>
                <div class="stat-card">
                    <div class="stat-label">개인사업자</div>
                    <div class="stat-value">{stats.get('total_indiv', 0):,}개</div>
                </div>
                <div class="stat-card {'warning' if stats.get('closure_ratio', 0) > 10 else ''}">
                    <div class="stat-label">폐업률</div>
                    <div class="stat-value">{stats.get('closure_ratio', 0)}%</div>
                    <div class="stat-sub">{stats.get('total_closed', 0):,}개 폐업</div>
                </div>
            </div>

            <div class="summary-box">
                <h3>핵심 요약</h3>
                <ul>
        """
            for s in insights.get('summary', []):
                insights_html += f"<li>{s}</li>"

            insights_html += """
                </ul>
            </div>

            <div class="highlights-grid">
        """

            for h in insights.get('highlights', []):
                insights_html += f"""
                <div class="highlight-card">
                    <div class="highlight-num">{h['icon']}</div>
                    <div class="highlight-content">
                        <div class="highlight-title">{h['title']}</div>
                        <div class="highlight-text">{h['content']}</div>
                    </div>
                </div>
            """

            insights_html += "</div>"

            # 경고사항 (지역별)
            if insights.get('warnings'):
                insights_html += '<div class="warnings-box"><h3>주의사항</h3><ul>'
                for w in insights['warnings']:
                    insights_html += f"<li>{w}</li>"
                insights_html += "</ul></div>"

            insights_html += "</div>"

    # 기준년월 옵션 (자료유형별 그룹화)
    base_ym_options_by_type = {}
    for item in filter_opts['base_ym_list']:
        dt = item['data_type']
        if dt not in base_ym_options_by_type:
            base_ym_options_by_type[dt] = []
        base_ym_options_by_type[dt].append(item)

    # 현재 선택된 자료유형의 기준년월 옵션
    current_base_ym_options = base_ym_options_by_type.get(selected_data_type, [])
    base_ym_options_html = ''.join([
        f'<option value="{item["base_ym"]}" {"selected" if item["base_ym"] == selected_base_ym else ""}>'
        f'{item["base_ym1"]}</option>'
        for item in current_base_ym_options
    ])

    # 자료유형 옵션
    data_type_options = ''.join([
        f'<option value="{dt}" {"selected" if dt == selected_data_type else ""}>{dt}</option>'
        for dt in filter_opts['data_type_list']
    ])

    # 시도 옵션 (sido_list가 시도명 리스트)
    sido_options = ''.join([
        f'<option value="{sido}" {"selected" if sido == selected_sido_nm else ""}>{sido}</option>'
        for sido in filter_opts['sido_list']
    ])

    # 기준년월 데이터를 JavaScript용으로 변환
    base_ym_data_js = {}
    for dt, items in base_ym_options_by_type.items():
        base_ym_data_js[dt] = [{'base_ym': i['base_ym'], 'base_ym1': i['base_ym1']} for i in items]

    import json
    base_ym_json = json.dumps(base_ym_data_js, ensure_ascii=False)

    # 테이블 데이터 생성
    if df_result is not None and len(df_result) > 0:
        # 테이블 헤더
        headers = ''.join([f'<th>{col}</th>' for col in df_result.columns])

        # 테이블 행
        rows = ''
        for _, row in df_result.iterrows():
            cells = ''.join([
                f'<td>{format_value(val)}</td>' for val in row
            ])
            rows += f'<tr>{cells}</tr>'

        table_html = f"""
        <table id="resultTable" class="display" style="width:100%">
            <thead><tr>{headers}</tr></thead>
            <tbody>{rows}</tbody>
        </table>
        """
        record_count = len(df_result)
    else:
        table_html = '<p class="no-data">조회된 데이터가 없습니다.</p>'
        record_count = 0

    html = f"""
<!DOCTYPE html>
<html lang="ko">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>기업체현황 분석 조회</title>
    <!-- DataTables CSS -->
    <link rel="stylesheet" href="https://cdn.datatables.net/1.13.7/css/jquery.dataTables.min.css">
    <link rel="stylesheet" href="https://cdn.datatables.net/buttons/2.4.2/css/buttons.dataTables.min.css">
    <style>
        * {{
            box-sizing: border-box;
            margin: 0;
            padding: 0;
        }}
        body {{
            font-family: 'Segoe UI', 'Malgun Gothic', sans-serif;
            background: linear-gradient(135deg, #3498db 0%, #2c3e50 100%);
            padding: 20px;
            min-height: 100vh;
        }}
        .container {{
            max-width: 1600px;
            margin: 0 auto;
            background: white;
            border-radius: 15px;
            padding: 30px;
            box-shadow: 0 10px 40px rgba(0,0,0,0.15);
        }}
        h1 {{
            color: #2c3e50;
            text-align: center;
            margin-bottom: 10px;
            font-size: 1.8rem;
        }}
        .subtitle {{
            text-align: center;
            color: #7f8c8d;
            margin-bottom: 25px;
            font-size: 0.95rem;
        }}
        .filter-container {{
            background: #f8f9fa;
            padding: 20px;
            border-radius: 12px;
            margin-bottom: 25px;
            border: 1px solid #e9ecef;
        }}
        .filter-row {{
            display: grid;
            grid-template-columns: repeat(auto-fit, minmax(180px, 1fr));
            gap: 15px;
            margin-bottom: 15px;
        }}
        .filter-group {{
            display: flex;
            flex-direction: column;
        }}
        .filter-group label {{
            font-size: 0.85rem;
            color: #2c3e50;
            margin-bottom: 5px;
            font-weight: 600;
        }}
        .filter-group select {{
            padding: 10px 12px;
            border: 1px solid #ced4da;
            border-radius: 6px;
            font-size: 0.95rem;
            background: white;
            cursor: pointer;
            transition: border-color 0.2s;
        }}
        .filter-group select:focus {{
            outline: none;
            border-color: #3498db;
            box-shadow: 0 0 0 3px rgba(52, 152, 219, 0.15);
        }}
        .radio-group {{
            display: flex;
            gap: 20px;
            flex-wrap: wrap;
            align-items: center;
            padding: 10px 0;
        }}
        .radio-label {{
            display: flex;
            align-items: center;
            cursor: pointer;
            font-size: 0.95rem;
            color: #2c3e50;
            font-weight: 500;
            padding: 8px 16px;
            border-radius: 20px;
            border: 2px solid #e9ecef;
            transition: all 0.2s;
        }}
        .radio-label:hover {{
            border-color: #3498db;
            background: #f0f7ff;
        }}
        .radio-label input {{
            margin-right: 8px;
            cursor: pointer;
            width: 16px;
            height: 16px;
        }}
        .radio-label input:checked + span {{
            color: #3498db;
            font-weight: 600;
        }}
        .radio-label:has(input:checked) {{
            border-color: #3498db;
            background: #e8f4fd;
        }}
        .radio-label.disabled {{
            opacity: 0.4;
            cursor: not-allowed;
            pointer-events: none;
            border-color: #ddd;
            background: #f5f5f5;
        }}
        .radio-label.disabled input {{
            cursor: not-allowed;
        }}
        .restriction-note {{
            font-size: 0.75rem;
            color: #999;
            margin-left: 4px;
        }}
        .radio-label.disabled .restriction-note {{
            color: #e74c3c;
        }}
        .btn-row {{
            display: flex;
            gap: 10px;
            justify-content: center;
            flex-wrap: wrap;
            margin-top: 15px;
        }}
        .filter-btn {{
            background: linear-gradient(135deg, #3498db, #2980b9);
            color: white;
            padding: 12px 35px;
            border: none;
            border-radius: 8px;
            font-size: 1rem;
            cursor: pointer;
            font-weight: 600;
            transition: all 0.2s;
        }}
        .filter-btn:hover {{
            transform: translateY(-2px);
            box-shadow: 0 4px 12px rgba(52, 152, 219, 0.4);
        }}
        .export-btn {{
            background: linear-gradient(135deg, #27ae60, #2ecc71);
            color: white;
            padding: 12px 25px;
            border: none;
            border-radius: 8px;
            font-size: 0.95rem;
            cursor: pointer;
            font-weight: 600;
            transition: all 0.2s;
        }}
        .export-btn:hover {{
            transform: translateY(-2px);
            box-shadow: 0 4px 12px rgba(39, 174, 96, 0.4);
        }}
        .export-btn.csv {{
            background: linear-gradient(135deg, #e67e22, #f39c12);
        }}
        .export-btn.csv:hover {{
            box-shadow: 0 4px 12px rgba(230, 126, 34, 0.4);
        }}
        .result-info {{
            display: flex;
            justify-content: space-between;
            align-items: center;
            margin-bottom: 15px;
            padding: 10px 15px;
            background: #e8f4fd;
            border-radius: 8px;
            flex-wrap: wrap;
            gap: 10px;
        }}
        .result-info .info-text {{
            color: #2c3e50;
            font-weight: 500;
        }}
        .result-info .record-count {{
            background: #3498db;
            color: white;
            padding: 5px 15px;
            border-radius: 20px;
            font-weight: 600;
        }}
        .table-container {{
            overflow-x: auto;
            margin-top: 15px;
        }}
        #resultTable {{
            width: 100%;
            border-collapse: collapse;
        }}
        #resultTable thead th {{
            background: #3498db;
            color: white;
            padding: 12px 15px;
            text-align: left;
            font-weight: 600;
            white-space: nowrap;
            position: sticky;
            top: 0;
        }}
        #resultTable tbody td {{
            padding: 10px 15px;
            border-bottom: 1px solid #e9ecef;
        }}
        #resultTable tbody tr:hover {{
            background: #f8f9fa;
        }}
        #resultTable tbody tr:nth-child(even) {{
            background: #fafbfc;
        }}
        .no-data {{
            text-align: center;
            padding: 50px;
            color: #7f8c8d;
            font-size: 1.1rem;
        }}
        /* DataTables 스타일 오버라이드 */
        .dataTables_wrapper .dataTables_length,
        .dataTables_wrapper .dataTables_filter {{
            margin-bottom: 15px;
        }}
        .dataTables_wrapper .dataTables_info,
        .dataTables_wrapper .dataTables_paginate {{
            margin-top: 15px;
        }}
        table.dataTable thead th {{
            cursor: pointer;
            position: relative;
            padding-right: 26px !important;
        }}
        /* 정렬 아이콘 스타일 개선 */
        table.dataTable thead th.sorting:after,
        table.dataTable thead th.sorting_asc:after,
        table.dataTable thead th.sorting_desc:after {{
            position: absolute;
            right: 8px;
            top: 50%;
            transform: translateY(-50%);
            font-size: 14px;
            opacity: 1 !important;
        }}
        table.dataTable thead th.sorting:after {{
            content: "\\25B2\\25BC" !important;
            color: rgba(255,255,255,0.4);
            font-size: 10px;
            line-height: 1;
        }}
        table.dataTable thead th.sorting_asc:after {{
            content: "\\25B2" !important;
            color: #ffeb3b;
            font-size: 14px;
        }}
        table.dataTable thead th.sorting_desc:after {{
            content: "\\25BC" !important;
            color: #ffeb3b;
            font-size: 14px;
        }}
        table.dataTable thead th:hover {{
            background: #2980b9;
        }}
        /* 시도 선택 표시/숨김 */
        .sido-select-group {{
            display: none;
        }}
        .sido-select-group.visible {{
            display: flex;
        }}
        /* 하위시군구 구분 옵션 표시/숨김 */
        .sigun-detail-group {{
            display: none;
        }}
        .sigun-detail-group.visible {{
            display: flex;
            flex-direction: column;
        }}
        /* 차트 스타일 */
        .chart-section {{
            margin-top: 30px;
            padding-top: 25px;
            border-top: 2px solid #e9ecef;
        }}
        .chart-section h2 {{
            color: #2c3e50;
            font-size: 1.3rem;
            margin-bottom: 20px;
            text-align: center;
        }}
        .chart-row {{
            display: grid;
            grid-template-columns: repeat(auto-fit, minmax(450px, 1fr));
            gap: 25px;
            margin-bottom: 25px;
        }}
        .chart-card {{
            background: #fafbfc;
            border-radius: 12px;
            padding: 15px;
            box-shadow: 0 2px 8px rgba(0,0,0,0.05);
            border: 1px solid #e9ecef;
        }}
        .chart-container {{
            width: 100%;
            min-height: 400px;
        }}
        /* 인사이트 섹션 스타일 */
        .insights-section {{
            background: linear-gradient(135deg, #f8f9fa 0%, #e9ecef 100%);
            border-radius: 15px;
            padding: 25px;
            margin-bottom: 25px;
            border: 1px solid #dee2e6;
        }}
        .insights-section h2 {{
            color: #2c3e50;
            font-size: 1.4rem;
            margin-bottom: 20px;
            padding-bottom: 10px;
            border-bottom: 2px solid #3498db;
        }}
        .stats-cards {{
            display: grid;
            grid-template-columns: repeat(auto-fit, minmax(180px, 1fr));
            gap: 15px;
            margin-bottom: 25px;
        }}
        .stat-card {{
            background: white;
            border-radius: 12px;
            padding: 20px;
            text-align: center;
            box-shadow: 0 3px 10px rgba(0,0,0,0.08);
            border-left: 4px solid #3498db;
            transition: transform 0.2s;
        }}
        .stat-card:hover {{
            transform: translateY(-3px);
        }}
        .stat-card.warning {{
            border-left-color: #e74c3c;
            background: #fff5f5;
        }}
        .stat-label {{
            font-size: 0.85rem;
            color: #7f8c8d;
            margin-bottom: 8px;
            font-weight: 600;
        }}
        .stat-value {{
            font-size: 1.8rem;
            font-weight: 700;
            color: #2c3e50;
        }}
        .stat-sub {{
            font-size: 0.8rem;
            color: #95a5a6;
            margin-top: 5px;
        }}
        .summary-box {{
            background: white;
            border-radius: 10px;
            padding: 20px;
            margin-bottom: 20px;
            border-left: 4px solid #27ae60;
        }}
        .summary-box h3 {{
            color: #27ae60;
            font-size: 1.1rem;
            margin-bottom: 12px;
        }}
        .summary-box ul {{
            list-style: none;
            padding: 0;
        }}
        .summary-box li {{
            padding: 8px 0;
            padding-left: 20px;
            position: relative;
            color: #2c3e50;
            font-size: 0.95rem;
            line-height: 1.5;
        }}
        .summary-box li:before {{
            content: "\\2022";
            color: #27ae60;
            font-weight: bold;
            position: absolute;
            left: 0;
        }}
        .highlights-grid {{
            display: grid;
            grid-template-columns: repeat(auto-fit, minmax(280px, 1fr));
            gap: 15px;
            margin-bottom: 20px;
        }}
        .highlight-card {{
            background: white;
            border-radius: 10px;
            padding: 18px;
            display: flex;
            align-items: flex-start;
            gap: 15px;
            box-shadow: 0 2px 8px rgba(0,0,0,0.06);
        }}
        .highlight-num {{
            background: linear-gradient(135deg, #3498db, #2980b9);
            color: white;
            width: 36px;
            height: 36px;
            border-radius: 50%;
            display: flex;
            align-items: center;
            justify-content: center;
            font-weight: 700;
            font-size: 1.1rem;
            flex-shrink: 0;
        }}
        .highlight-content {{
            flex: 1;
        }}
        .highlight-title {{
            font-weight: 600;
            color: #2c3e50;
            margin-bottom: 5px;
            font-size: 0.95rem;
        }}
        .highlight-text {{
            color: #5d6d7e;
            font-size: 0.9rem;
            line-height: 1.4;
        }}
        .warnings-box {{
            background: #fff5f5;
            border-radius: 10px;
            padding: 15px 20px;
            border-left: 4px solid #e74c3c;
        }}
        .warnings-box h3 {{
            color: #e74c3c;
            font-size: 1rem;
            margin-bottom: 10px;
        }}
        .warnings-box ul {{
            list-style: none;
            padding: 0;
            margin: 0;
        }}
        .warnings-box li {{
            color: #c0392b;
            font-size: 0.9rem;
            padding: 5px 0;
        }}
        .ppt-btn {{
            background: linear-gradient(135deg, #e67e22, #d35400);
            color: white;
            padding: 12px 25px;
            border: none;
            border-radius: 8px;
            font-size: 0.95rem;
            cursor: pointer;
            font-weight: 600;
            transition: all 0.2s;
        }}
        .ppt-btn:hover {{
            transform: translateY(-2px);
            box-shadow: 0 4px 12px rgba(230, 126, 34, 0.4);
        }}
        .ppt-btn:disabled {{
            background: #bdc3c7;
            cursor: not-allowed;
            transform: none;
        }}
        @media (max-width: 768px) {{
            .chart-row {{
                grid-template-columns: 1fr;
            }}
            .stats-cards {{
                grid-template-columns: repeat(2, 1fr);
            }}
            .highlights-grid {{
                grid-template-columns: 1fr;
            }}
        }}
    </style>
</head>
<body>
    <div class="container">
        <h1>기업체현황 분석 조회</h1>
        <p class="subtitle">지역별(권역/시도/시군구) 및 항목별(기업규모/산업분류/조직형태 등) 기업체 현황을 조회합니다</p>

        <div class="filter-container">
            <form id="filterForm" method="GET">
                <div class="filter-row">
                    <div class="filter-group">
                        <label>자료유형</label>
                        <select name="data_type" id="dataTypeSelect">
                            {data_type_options}
                        </select>
                    </div>
                    <div class="filter-group">
                        <label>기준년월</label>
                        <select name="base_ym" id="baseYmSelect">
                            {base_ym_options_html}
                        </select>
                    </div>
                    <div class="filter-group sido-select-group {'visible' if selected_region_type == '시군구별' else ''}" id="sidoSelectGroup">
                        <label>시도 선택</label>
                        <select name="sido_nm" id="sidoSelect">
                            {sido_options}
                        </select>
                    </div>
                </div>

                <div class="filter-group">
                    <label>지역 구분</label>
                    <div class="radio-group">
                        <label class="radio-label">
                            <input type="radio" name="region_type" value="권역별"
                                   {'checked' if selected_region_type == '권역별' else ''}>
                            <span>권역별</span>
                        </label>
                        <label class="radio-label">
                            <input type="radio" name="region_type" value="시도별"
                                   {'checked' if selected_region_type == '시도별' else ''}>
                            <span>시도별</span>
                        </label>
                        <label class="radio-label">
                            <input type="radio" name="region_type" value="시군구별"
                                   {'checked' if selected_region_type == '시군구별' else ''}>
                            <span>시군구별</span>
                        </label>
                    </div>
                </div>

                <div class="filter-group">
                    <label>항목 구분 (선택사항)</label>
                    <div class="radio-group" id="itemAnalysisGroup">
                        <label class="radio-label" data-types="all">
                            <input type="radio" name="item_type" value=""
                                   {'checked' if not selected_item_type else ''}>
                            <span>기본 현황</span>
                        </label>
                        <label class="radio-label" data-types="annual">
                            <input type="radio" name="item_type" value="기업규모별"
                                   {'checked' if selected_item_type == '기업규모별' else ''}>
                            <span>기업규모별</span>
                            <small class="restriction-note">(연간만)</small>
                        </label>
                        <label class="radio-label" data-types="all">
                            <input type="radio" name="item_type" value="산업분류별"
                                   {'checked' if selected_item_type == '산업분류별' else ''}>
                            <span>산업분류별</span>
                        </label>
                        <label class="radio-label" data-types="all">
                            <input type="radio" name="item_type" value="조직형태별"
                                   {'checked' if selected_item_type == '조직형태별' else ''}>
                            <span>조직형태별</span>
                        </label>
                        <label class="radio-label" data-types="all">
                            <input type="radio" name="item_type" value="대표자성별별"
                                   {'checked' if selected_item_type == '대표자성별별' else ''}>
                            <span>대표자성별</span>
                        </label>
                        <label class="radio-label" data-types="annual,monthly">
                            <input type="radio" name="item_type" value="연령그룹별"
                                   {'checked' if selected_item_type == '연령그룹별' else ''}>
                            <span>연령그룹별</span>
                            <small class="restriction-note">(분기 제외)</small>
                        </label>
                        <label class="radio-label" data-types="all">
                            <input type="radio" name="item_type" value="폐업여부별"
                                   {'checked' if selected_item_type == '폐업여부별' else ''}>
                            <span>폐업여부별</span>
                        </label>
                    </div>
                </div>

                <!-- 하위시군구 구분 옵션 (시군구별 선택 시만 표시) -->
                <div class="filter-group sigun-detail-group {'visible' if selected_region_type == '시군구별' else ''}" id="sigunDetailGroup">
                    <label>하위시군구 구분</label>
                    <div class="radio-group">
                        <label class="radio-label">
                            <input type="radio" name="sigun_detail" value="false"
                                   {'checked' if not selected_sigun_detail else ''}>
                            <span>구분없이 (합계)</span>
                        </label>
                        <label class="radio-label">
                            <input type="radio" name="sigun_detail" value="true"
                                   {'checked' if selected_sigun_detail else ''}>
                            <span>상세 표시</span>
                        </label>
                    </div>
                </div>

                <div class="btn-row">
                    <button type="submit" class="filter-btn">조회</button>
                    <button type="button" class="export-btn" onclick="exportData('excel')">Excel 다운로드</button>
                    <button type="button" class="export-btn csv" onclick="exportData('csv')">CSV 다운로드</button>
                    <button type="button" class="ppt-btn" onclick="savePPT()" {'disabled' if not PPT_AVAILABLE else ''}>
                        {'PPT 다운로드' if PPT_AVAILABLE else 'PPT (미설치)'}
                    </button>
                </div>
            </form>

            <!-- PPT 저장 폼 (숨김) -->
            <form id="pptForm" method="POST" style="display:none;">
                <input type="hidden" name="action" value="save_ppt">
                <input type="hidden" name="base_ym" value="{selected_base_ym}">
                <input type="hidden" name="data_type" value="{selected_data_type}">
                <input type="hidden" name="analysis_type" value="{selected_region_type}">
                <input type="hidden" name="item_type" value="{selected_item_type if selected_item_type else ''}">
                <input type="hidden" name="sido_nm" value="{selected_sido_nm}">
                <input type="hidden" name="sigun_detail" value="{'true' if selected_sigun_detail else 'false'}">
            </form>
        </div>

        <!-- 인사이트 섹션 -->
        {insights_html}

        <div class="result-info">
            <span class="info-text">
                <strong>{selected_data_type}</strong> |
                <strong>{selected_base_ym}</strong> |
                <strong>{selected_analysis_type}</strong>
                {f' | <strong>{selected_sido_nm}</strong>' if selected_region_type == '시군구별' and selected_sido_nm else ''}
            </span>
            <span class="record-count">총 {record_count:,}건</span>
        </div>

        <div class="table-container">
            {table_html}
        </div>

        <!-- 차트 섹션 -->
        {f'''
        <div class="chart-section">
            <h2>시각화 분석</h2>
            <div class="chart-row">
                <div class="chart-card">
                    <div id="barChart" class="chart-container"></div>
                </div>
                <div class="chart-card">
                    <div id="pieChart" class="chart-container"></div>
                </div>
            </div>
            <div class="chart-row">
                <div class="chart-card" style="max-width: 800px; margin: 0 auto;">
                    <div id="closureChart" class="chart-container"></div>
                </div>
            </div>
        </div>
        ''' if bar_chart_json else ''}
    </div>

    <!-- Plotly JS -->
    <script src="https://cdn.plot.ly/plotly-2.27.0.min.js"></script>
    <!-- jQuery -->
    <script src="https://code.jquery.com/jquery-3.7.1.min.js"></script>
    <!-- DataTables JS -->
    <script src="https://cdn.datatables.net/1.13.7/js/jquery.dataTables.min.js"></script>

    <script>
        // 기준년월 데이터
        const baseYmData = {base_ym_json};

        // 분석유형 활성화/비활성화 처리 함수
        function updateAnalysisTypeAvailability(dataType) {{
            const itemAnalysisGroup = document.getElementById('itemAnalysisGroup');
            if (!itemAnalysisGroup) return;

            const labels = itemAnalysisGroup.querySelectorAll('.radio-label[data-types]');
            labels.forEach(label => {{
                const allowedTypes = label.getAttribute('data-types');
                const input = label.querySelector('input');

                // 허용된 자료유형인지 확인
                let isAllowed = false;
                if (allowedTypes === 'all') {{
                    isAllowed = true;
                }} else {{
                    const typeList = allowedTypes.split(',');
                    isAllowed = typeList.includes(dataType);
                }}

                if (isAllowed) {{
                    label.classList.remove('disabled');
                    input.disabled = false;
                }} else {{
                    label.classList.add('disabled');
                    input.disabled = true;
                    // 비활성화된 항목이 선택되어 있으면 권역별로 변경
                    if (input.checked) {{
                        const regionRadio = document.querySelector('input[name="analysis_type"][value="권역별"]');
                        if (regionRadio) {{
                            regionRadio.checked = true;
                        }}
                    }}
                }}
            }});
        }}

        // 자료유형 변경 시 기준년월 옵션 업데이트 및 분석유형 활성화/비활성화
        document.getElementById('dataTypeSelect').addEventListener('change', function() {{
            const dataType = this.value;
            const baseYmSelect = document.getElementById('baseYmSelect');
            const options = baseYmData[dataType] || [];

            baseYmSelect.innerHTML = '';
            options.forEach(item => {{
                const opt = document.createElement('option');
                opt.value = item.base_ym;
                opt.textContent = item.base_ym1;
                baseYmSelect.appendChild(opt);
            }});

            // 분석유형 활성화/비활성화 업데이트
            updateAnalysisTypeAvailability(dataType);
        }});

        // 페이지 로드 시 초기 상태 설정
        document.addEventListener('DOMContentLoaded', function() {{
            const dataTypeSelect = document.getElementById('dataTypeSelect');
            if (dataTypeSelect) {{
                updateAnalysisTypeAvailability(dataTypeSelect.value);
            }}
        }});

        // 지역구분 변경 시 시도 선택 및 하위시군구 구분 옵션 표시/숨김
        document.querySelectorAll('input[name="region_type"]').forEach(radio => {{
            radio.addEventListener('change', function() {{
                const sidoGroup = document.getElementById('sidoSelectGroup');
                const sigunDetailGroup = document.getElementById('sigunDetailGroup');
                if (this.value === '시군구별') {{
                    sidoGroup.classList.add('visible');
                    sigunDetailGroup.classList.add('visible');
                }} else {{
                    sidoGroup.classList.remove('visible');
                    sigunDetailGroup.classList.remove('visible');
                }}
            }});
        }});

        // DataTables 초기화
        $(document).ready(function() {{
            if ($('#resultTable').length) {{
                $('#resultTable').DataTable({{
                    language: {{
                        search: "검색:",
                        lengthMenu: "_MENU_ 건씩 보기",
                        info: "_START_ - _END_ / _TOTAL_ 건",
                        infoEmpty: "데이터 없음",
                        infoFiltered: "(전체 _MAX_ 건 중)",
                        paginate: {{
                            first: "처음",
                            last: "끝",
                            next: "다음",
                            previous: "이전"
                        }},
                        emptyTable: "데이터가 없습니다",
                        zeroRecords: "검색 결과가 없습니다"
                    }},
                    pageLength: 25,
                    lengthMenu: [[10, 25, 50, 100, -1], [10, 25, 50, 100, "전체"]],
                    order: [],  // 초기 정렬 없음
                    scrollX: true
                }});
            }}
        }});

        // 데이터 내보내기
        function exportData(format) {{
            const form = document.getElementById('filterForm');
            const formData = new FormData(form);
            formData.append('export', format);

            const params = new URLSearchParams(formData);
            window.location.href = window.location.pathname + '?' + params.toString();
        }}

        // PPT 저장
        function savePPT() {{
            document.getElementById('pptForm').submit();
        }}

        // 차트 렌더링
        {f'''
        (function() {{
            // 막대차트
            const barChartData = {bar_chart_json};
            if (barChartData && document.getElementById('barChart')) {{
                Plotly.newPlot('barChart', barChartData.data, barChartData.layout, {{responsive: true}});
            }}

            // 파이차트
            const pieChartData = {pie_chart_json if pie_chart_json else 'null'};
            if (pieChartData && document.getElementById('pieChart')) {{
                Plotly.newPlot('pieChart', pieChartData.data, pieChartData.layout, {{responsive: true}});
            }}

            // 폐업률 차트
            const closureChartData = {closure_chart_json if closure_chart_json else 'null'};
            if (closureChartData && document.getElementById('closureChart')) {{
                Plotly.newPlot('closureChart', closureChartData.data, closureChartData.layout, {{responsive: true}});
            }}
        }})();
        ''' if bar_chart_json else ''}
    </script>
</body>
</html>
    """

    return html


def format_value(val):
    """값 포맷팅"""
    if pd.isna(val):
        return '-'
    elif isinstance(val, float):
        if val == int(val):
            return f'{int(val):,}'
        return f'{val:,.2f}'
    elif isinstance(val, int):
        return f'{val:,}'
    return str(val)


def render(request_args=None):
    """
    Flask에서 호출되는 렌더링 함수

    Args:
        request_args (dict, optional): Flask request.args

    Returns:
        str: 렌더링된 HTML 또는 Flask Response (내보내기/PPT 저장 시)
    """
    logger.info("=== 기업체현황 분석 조회 ===")

    # 필터 옵션 가져오기
    filter_opts = get_filter_options()

    # POST 요청 처리 (PPT 저장)
    try:
        if request.method == 'POST':
            action = request.form.get('action')
            logger.info(f"POST 요청 - action: {action}")

            if action == 'save_ppt':
                # 폼 데이터에서 파라미터 추출
                form_base_ym = request.form.get('base_ym', '')
                form_data_type = request.form.get('data_type', '')
                form_analysis_type = request.form.get('analysis_type', '시도별')
                form_item_type = request.form.get('item_type', '')
                form_sido_nm = request.form.get('sido_nm', '')
                form_sigun_detail = request.form.get('sigun_detail', 'false') == 'true'

                # 데이터 조회
                df_for_ppt = get_analysis_data(
                    form_base_ym,
                    form_data_type,
                    form_analysis_type,
                    form_sido_nm if form_analysis_type == '시군구별' else None,
                    form_sigun_detail,
                    form_item_type if form_item_type else None
                )

                # 인사이트 생성
                insights_for_ppt = generate_insights(df_for_ppt, form_analysis_type, form_base_ym, form_data_type, form_item_type if form_item_type else None)

                return handle_ppt_save(request.form, df_for_ppt, insights_for_ppt, filter_opts)
    except RuntimeError as e:
        # request context 외부에서 호출된 경우 무시
        logger.debug(f"RuntimeError (무시됨): {e}")
    except Exception as e:
        logger.error(f"POST 처리 중 예외: {e}")
        import traceback
        logger.error(traceback.format_exc())

    # 요청 파라미터 파싱
    request_args = request_args or {}

    # 기본값 설정
    default_data_type = filter_opts['data_type_list'][0] if filter_opts['data_type_list'] else '연간'
    default_base_ym = filter_opts['base_ym_list'][0]['base_ym'] if filter_opts['base_ym_list'] else ''
    default_sido_nm = filter_opts['sido_list'][0] if filter_opts['sido_list'] else ''

    selected_data_type = request_args.get('data_type', default_data_type)
    selected_base_ym = request_args.get('base_ym', default_base_ym)
    selected_region_type = request_args.get('region_type', '시도별')  # 지역 구분
    selected_item_type = request_args.get('item_type', '')  # 항목 구분 (선택사항)
    selected_sido_nm = request_args.get('sido_nm', default_sido_nm)
    # 하위시군구 구분: 기본값은 False (구분없이 합계)
    selected_sigun_detail = request_args.get('sigun_detail', 'false') == 'true'
    export_format = request_args.get('export', None)

    # 기존 호환성: analysis_type 파라미터도 처리
    legacy_analysis_type = request_args.get('analysis_type', None)
    if legacy_analysis_type:
        if legacy_analysis_type in ['권역별', '시도별', '시군구별']:
            selected_region_type = legacy_analysis_type
        else:
            selected_item_type = legacy_analysis_type

    logger.info(f"파라미터: data_type={selected_data_type}, base_ym={selected_base_ym}, "
                f"region_type={selected_region_type}, item_type={selected_item_type}, "
                f"sido_nm={selected_sido_nm}, sigun_detail={selected_sigun_detail}")

    # 데이터 조회
    df_result = get_analysis_data(
        selected_base_ym,
        selected_data_type,
        selected_region_type,
        selected_sido_nm if selected_region_type == '시군구별' else None,
        selected_sigun_detail,
        selected_item_type  # 항목 구분 추가
    )

    logger.info(f"조회 결과: {len(df_result)}건")

    # 분석유형 표시용 문자열 생성
    analysis_label = selected_region_type
    if selected_item_type:
        analysis_label = f"{selected_region_type}_{selected_item_type}"

    # 내보내기 요청 처리
    if export_format in ['excel', 'csv']:
        filename_prefix = f"기업체현황_{analysis_label}_{selected_base_ym}"
        return handle_export(df_result, export_format, filename_prefix)

    # 인사이트 생성
    insights = generate_insights(df_result, selected_region_type, selected_base_ym, selected_data_type, selected_item_type)

    # HTML 생성
    html = generate_html(
        filter_opts,
        selected_base_ym,
        selected_data_type,
        selected_region_type,
        selected_item_type,
        selected_sido_nm,
        selected_sigun_detail,
        df_result,
        insights
    )

    return html
