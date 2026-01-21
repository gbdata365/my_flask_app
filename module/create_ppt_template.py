# -*- coding: utf-8 -*-
"""
PPT 템플릿 생성 스크립트 (PowerPoint Template Generator)
======================================================

기본 PPT 템플릿 파일을 프로그래밍 방식으로 생성하는 스크립트입니다.
생성된 템플릿을 PowerPoint에서 열어 마스터 슬라이드를 수정하여
조직의 브랜딩에 맞게 커스터마이징할 수 있습니다.

주요 기능:
---------
1. 8종류의 샘플 슬라이드 레이아웃 생성
2. 일관된 색상 테마 및 스타일 적용
3. 16:9 와이드스크린 비율 지원

생성되는 슬라이드 종류:
-------------------
1. 제목 슬라이드 (Title Slide)
2. 목차 슬라이드 (Table of Contents)
3. 내용 슬라이드 (Content - 제목 + 본문)
4. 지표 카드 슬라이드 (Metrics Cards - 2x2 그리드)
5. 차트 슬라이드 (Chart Placeholder)
6. 표 슬라이드 (Table)
7. 인사이트 슬라이드 (Insights)
8. 마무리 슬라이드 (Thank You)

사용법:
------
    # 방법 1: 직접 실행
    $ python module/create_ppt_template.py

    # 방법 2: 모듈 임포트
    from module.create_ppt_template import create_template
    prs = create_template()
    prs.save('my_template.pptx')

출력 위치:
---------
    templates/report_template.pptx

템플릿 수정 방법:
--------------
    1. PowerPoint에서 생성된 템플릿 파일 열기
    2. 보기 → 슬라이드 마스터 클릭
    3. 각 레이아웃의 색상, 폰트, 배경 등 수정
    4. 슬라이드 마스터 닫기 → 파일 저장

색상 테마:
--------
    - 기본 배경: #1243A6 (진한 파랑)
    - 카드 색상: #667EEA (파랑), #764BA2 (보라), #2ECC71 (초록), #3498DB (하늘)
    - 텍스트: #FFFFFF (흰색), #323232 (진한 회색)
    - 테두리/구분선: #C8C8C8 (연한 회색)

의존성:
------
    - python-pptx: PPT 파일 생성 라이브러리
    - pathlib: 크로스 플랫폼 경로 처리

Author: Claude AI Agent
Created: 2025-01-15
License: MIT
"""

# =============================================================================
# 표준 라이브러리 임포트
# =============================================================================
from pathlib import Path

# =============================================================================
# 서드파티 라이브러리 임포트 (python-pptx)
# =============================================================================
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml


# =============================================================================
# 색상 상수 정의 (Color Constants)
# 일관된 디자인을 위한 색상 팔레트
# =============================================================================

# 기본 브랜드 색상
PRIMARY_DARK = RGBColor(18, 67, 166)      # #1243A6 - 진한 파랑 (헤더 배경)
PRIMARY = RGBColor(102, 126, 234)          # #667EEA - 밝은 파랑 (카드 배경)

# 카드 색상 (순환 사용)
CARD_COLORS = [
    RGBColor(102, 126, 234),  # #667EEA - 파랑
    RGBColor(118, 75, 162),   # #764BA2 - 보라
    RGBColor(46, 204, 113),   # #2ECC71 - 초록
    RGBColor(52, 152, 219),   # #3498DB - 하늘
]

# 텍스트 색상
TEXT_WHITE = RGBColor(255, 255, 255)       # #FFFFFF - 흰색
TEXT_LIGHT = RGBColor(200, 220, 255)       # #C8DCFF - 연한 흰색
TEXT_GRAY = RGBColor(50, 50, 50)           # #323232 - 진한 회색
TEXT_MUTED = RGBColor(150, 150, 150)       # #969696 - 연한 회색

# 배경 색상
BG_LIGHT = RGBColor(248, 249, 250)         # #F8F9FA - 연한 회색 배경
BG_TABLE_HEADER = RGBColor(18, 67, 166)    # #1243A6 - 표 헤더
BG_TABLE_ROW = RGBColor(245, 247, 250)     # #F5F7FA - 표 짝수 행

# 테두리 색상
BORDER_LIGHT = RGBColor(200, 200, 200)     # #C8C8C8 - 연한 회색
BORDER_CARD = RGBColor(220, 220, 220)      # #DCDCDC - 카드 테두리


# =============================================================================
# 슬라이드 크기 상수
# =============================================================================
SLIDE_WIDTH = Inches(13.333)   # 16:9 와이드스크린
SLIDE_HEIGHT = Inches(7.5)


# =============================================================================
# 메인 함수: create_template
# =============================================================================

def create_template() -> Presentation:
    """
    기본 PPT 템플릿 생성

    8종류의 슬라이드 레이아웃이 포함된 PPT 템플릿을 생성합니다.
    각 슬라이드는 보고서에서 자주 사용되는 형식으로 구성됩니다.

    Returns:
        Presentation: python-pptx Presentation 객체

    생성되는 슬라이드:
    ---------------
    1. 제목 슬라이드 - 보고서 첫 페이지
    2. 목차 슬라이드 - 순서 안내
    3. 내용 슬라이드 - 일반 텍스트 내용
    4. 지표 카드 슬라이드 - KPI 표시 (2x2)
    5. 차트 슬라이드 - 시각화 영역
    6. 표 슬라이드 - 데이터 테이블
    7. 인사이트 슬라이드 - 분석 결과
    8. 마무리 슬라이드 - 감사 인사

    Example:
        >>> prs = create_template()
        >>> prs.save('my_template.pptx')
        >>>
        >>> # 슬라이드 수 확인
        >>> print(len(prs.slides))
        8
    """
    # -------------------------------------------------------------------------
    # 새 프레젠테이션 생성
    # -------------------------------------------------------------------------
    prs = Presentation()

    # 슬라이드 크기 설정 (16:9 와이드스크린)
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # -------------------------------------------------------------------------
    # 각 슬라이드 추가
    # -------------------------------------------------------------------------
    _add_title_slide(prs)           # 1. 제목 슬라이드
    _add_toc_slide(prs)             # 2. 목차 슬라이드
    _add_content_slide(prs)         # 3. 내용 슬라이드
    _add_metrics_slide(prs)         # 4. 지표 카드 슬라이드
    _add_chart_slide(prs)           # 5. 차트 슬라이드
    _add_table_slide(prs)           # 6. 표 슬라이드
    _add_insights_slide(prs)        # 7. 인사이트 슬라이드
    _add_thankyou_slide(prs)        # 8. 마무리 슬라이드

    return prs


# =============================================================================
# 슬라이드 생성 함수들 (Slide Creation Functions)
# =============================================================================

def _add_title_slide(prs: Presentation) -> None:
    """
    제목 슬라이드 추가 (1번 슬라이드)

    보고서의 첫 페이지로 사용되는 제목 슬라이드를 생성합니다.
    진한 파란색 배경에 중앙 정렬된 제목과 부제목을 배치합니다.

    레이아웃:
    --------
        ┌─────────────────────────────────────┐
        │                                     │
        │                                     │
        │     보고서 제목을 입력하세요          │  ← 44pt, 볼드, 흰색
        │                                     │
        │     부제목 | 날짜 | 작성자           │  ← 20pt, 연한 색
        │                                     │
        │                                     │
        │     기관명 또는 로고                 │  ← 하단 푸터
        └─────────────────────────────────────┘

    Args:
        prs (Presentation): 프레젠테이션 객체

    스타일:
        - 배경: #1243A6 (진한 파랑)
        - 제목: 44pt, 볼드, 흰색, 중앙 정렬
        - 부제목: 20pt, 연한 흰색, 중앙 정렬
        - 푸터: 14pt, 연한 흰색
    """
    # 빈 슬라이드 레이아웃 사용 (인덱스 6)
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # -------------------------------------------------------------------------
    # 배경색 설정 (전체 슬라이드를 덮는 사각형으로 구현)
    # -------------------------------------------------------------------------
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = PRIMARY_DARK
    bg_shape.line.fill.background()  # 테두리 없음

    # -------------------------------------------------------------------------
    # 제목 텍스트 박스
    # 위치: 화면 중앙
    # -------------------------------------------------------------------------
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5),      # 좌상단 위치
        Inches(12.333), Inches(1.5)    # 너비 x 높이
    )
    tf = title_box.text_frame
    tf.word_wrap = True  # 긴 제목 자동 줄바꿈

    p = tf.paragraphs[0]
    p.text = "보고서 제목을 입력하세요"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.alignment = PP_ALIGN.CENTER

    # -------------------------------------------------------------------------
    # 부제목 텍스트 박스
    # 위치: 제목 아래
    # -------------------------------------------------------------------------
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(4.2),
        Inches(12.333), Inches(0.8)
    )
    tf2 = subtitle_box.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "부제목 | 날짜 | 작성자"
    p2.font.size = Pt(20)
    p2.font.color.rgb = TEXT_LIGHT
    p2.alignment = PP_ALIGN.CENTER

    # -------------------------------------------------------------------------
    # 하단 푸터 영역 (기관명/로고)
    # 위치: 슬라이드 하단
    # -------------------------------------------------------------------------
    footer_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(6.5),
        Inches(12.333), Inches(0.5)
    )
    tf3 = footer_box.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "기관명 또는 로고"
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(180, 200, 230)  # 연한 흰색
    p3.alignment = PP_ALIGN.CENTER


def _add_toc_slide(prs: Presentation) -> None:
    """
    목차 슬라이드 추가 (2번 슬라이드)

    보고서의 순서를 안내하는 목차 슬라이드를 생성합니다.
    상단에 파란색 헤더 바와 함께 번호가 매겨진 항목을 나열합니다.

    레이아웃:
    --------
        ┌─────────────────────────────────────┐
        │ 목차                                │  ← 파란색 헤더 바
        ├─────────────────────────────────────┤
        │                                     │
        │   1. 개요                           │
        │                                     │
        │   2. 주요 지표                       │
        │                                     │
        │   3. 상세 분석                       │
        │                                     │
        │   4. 결론                           │
        │                                     │
        └─────────────────────────────────────┘

    Args:
        prs (Presentation): 프레젠테이션 객체

    스타일:
        - 헤더 바: #1243A6, 높이 1.2"
        - 제목: 28pt, 볼드, 흰색
        - 목차 항목: 24pt, 진한 회색, 세로 간격 1.2"
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # -------------------------------------------------------------------------
    # 상단 헤더 바 (파란색)
    # -------------------------------------------------------------------------
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_DARK
    header.line.fill.background()

    # -------------------------------------------------------------------------
    # 슬라이드 제목 ("목차")
    # -------------------------------------------------------------------------
    title = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.35),
        Inches(6), Inches(0.6)
    )
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "목차"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE

    # -------------------------------------------------------------------------
    # 목차 항목들
    # 세로로 간격을 두고 배치
    # -------------------------------------------------------------------------
    items = ["1. 개요", "2. 주요 지표", "3. 상세 분석", "4. 결론"]

    for i, item in enumerate(items):
        item_box = slide.shapes.add_textbox(
            Inches(1),                    # 왼쪽 들여쓰기
            Inches(1.8 + i * 1.2),        # 세로 위치 (1.2" 간격)
            Inches(10), Inches(0.8)
        )
        tf = item_box.text_frame
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(24)
        p.font.color.rgb = TEXT_GRAY


def _add_content_slide(prs: Presentation) -> None:
    """
    내용 슬라이드 추가 (3번 슬라이드)

    일반적인 텍스트 내용을 담는 슬라이드를 생성합니다.
    제목과 불릿 포인트 형식의 본문 영역으로 구성됩니다.

    레이아웃:
    --------
        ┌─────────────────────────────────────┐
        │ 슬라이드 제목                        │  ← 파란색 헤더 바
        ├─────────────────────────────────────┤
        │                                     │
        │   본문 내용을 입력하세요.             │
        │                                     │
        │   • 첫 번째 항목                     │
        │   • 두 번째 항목                     │
        │   • 세 번째 항목                     │
        │                                     │
        └─────────────────────────────────────┘

    Args:
        prs (Presentation): 프레젠테이션 객체

    스타일:
        - 헤더 바: #1243A6, 높이 1.2"
        - 제목: 28pt, 볼드, 흰색
        - 본문: 18pt, 진한 회색, 줄 간격 1.5
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 상단 헤더 바
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_DARK
    header.line.fill.background()

    # 슬라이드 제목
    title = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.35),
        Inches(10), Inches(0.6)
    )
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "슬라이드 제목"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE

    # -------------------------------------------------------------------------
    # 본문 영역
    # 불릿 포인트 형식의 예시 텍스트
    # -------------------------------------------------------------------------
    body_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.5),
        Inches(12.333), Inches(5.5)
    )
    tf = body_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "본문 내용을 입력하세요.\n\n• 첫 번째 항목\n• 두 번째 항목\n• 세 번째 항목"
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_GRAY
    p.line_spacing = 1.5  # 줄 간격


def _add_metrics_slide(prs: Presentation) -> None:
    """
    지표 카드 슬라이드 추가 (4번 슬라이드)

    KPI(핵심 성과 지표)를 2x2 그리드의 카드 형태로 표시하는 슬라이드입니다.
    대시보드의 주요 지표를 한눈에 보여줄 때 사용합니다.

    레이아웃:
    --------
        ┌─────────────────────────────────────┐
        │ 주요 지표                           │  ← 파란색 헤더 바
        ├─────────────────────────────────────┤
        │                                     │
        │   ┌─────────┐  ┌─────────┐         │
        │   │ 지표1   │  │ 지표2   │         │  ← 2x2 카드 그리드
        │   │ 123,456 │  │ 567,890 │         │
        │   │ 개      │  │ 명      │         │
        │   └─────────┘  └─────────┘         │
        │                                     │
        │   ┌─────────┐  ┌─────────┐         │
        │   │ 지표3   │  │ 지표4   │         │
        │   │ 1,523   │  │ +2.5%   │         │
        │   └─────────┘  └─────────┘         │
        │                                     │
        └─────────────────────────────────────┘

    Args:
        prs (Presentation): 프레젠테이션 객체

    스타일:
        - 카드 크기: 5.8" x 2.5"
        - 카드 색상: 파랑, 보라, 초록, 하늘 (순환)
        - 라벨: 16pt, 흰색
        - 값: 40pt, 볼드, 흰색
        - 단위: 14pt, 연한 흰색
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 헤더
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_DARK
    header.line.fill.background()

    title = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.35),
        Inches(10), Inches(0.6)
    )
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "주요 지표"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE

    # -------------------------------------------------------------------------
    # 4개 카드 배치 (2x2 그리드)
    # -------------------------------------------------------------------------
    card_data = [
        ("총 사업체수", "123,456", "개"),
        ("총 종사자수", "567,890", "명"),
        ("평균 HHI", "1,523", ""),
        ("증감률", "+2.5%", "전분기 대비"),
    ]

    # 카드 레이아웃 설정
    card_width = Inches(5.8)
    card_height = Inches(2.5)
    start_left = Inches(0.5)
    start_top = Inches(1.5)
    gap = Inches(0.4)

    for i, (label, value, unit) in enumerate(card_data):
        # 그리드 위치 계산
        row = i // 2  # 0 또는 1
        col = i % 2   # 0 또는 1

        left = start_left + col * (card_width + gap)
        top = start_top + row * (card_height + gap)

        # 카드 배경 (둥근 사각형)
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, card_width, card_height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_COLORS[i]
        card.line.fill.background()

        # 라벨 (지표명)
        label_box = slide.shapes.add_textbox(
            left + Inches(0.3), top + Inches(0.3),
            card_width - Inches(0.6), Inches(0.5)
        )
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_WHITE

        # 값 (큰 숫자)
        value_box = slide.shapes.add_textbox(
            left + Inches(0.3), top + Inches(0.9),
            card_width - Inches(0.6), Inches(1)
        )
        tf = value_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE

        # 단위
        unit_box = slide.shapes.add_textbox(
            left + Inches(0.3), top + Inches(1.9),
            card_width - Inches(0.6), Inches(0.4)
        )
        tf = unit_box.text_frame
        p = tf.paragraphs[0]
        p.text = unit
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(230, 230, 230)


def _add_chart_slide(prs: Presentation) -> None:
    """
    차트 슬라이드 추가 (5번 슬라이드)

    시각화(차트/그래프)를 배치할 플레이스홀더가 있는 슬라이드입니다.
    matplotlib 등으로 생성한 차트 이미지를 이 영역에 삽입합니다.

    레이아웃:
    --------
        ┌─────────────────────────────────────┐
        │ 차트 제목                           │  ← 파란색 헤더 바
        ├─────────────────────────────────────┤
        │                                     │
        │   ┌───────────────────────────┐     │
        │   │                           │     │
        │   │    차트 이미지가           │     │  ← 플레이스홀더 영역
        │   │    들어갈 영역             │     │
        │   │                           │     │
        │   └───────────────────────────┘     │
        │                                     │
        └─────────────────────────────────────┘

    Args:
        prs (Presentation): 프레젠테이션 객체

    스타일:
        - 차트 영역: 연한 회색 배경 (#F8F9FA)
        - 테두리: 연한 회색 (#C8C8C8)
        - 안내 텍스트: 18pt, 회색, 중앙 정렬
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 헤더
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_DARK
    header.line.fill.background()

    title = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.35),
        Inches(10), Inches(0.6)
    )
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "차트 제목"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE

    # -------------------------------------------------------------------------
    # 차트 플레이스홀더 영역
    # 실제 사용 시 이 영역에 차트 이미지 삽입
    # -------------------------------------------------------------------------
    chart_area = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(1.5),
        Inches(12.333), Inches(5.5)
    )
    chart_area.fill.solid()
    chart_area.fill.fore_color.rgb = BG_LIGHT
    chart_area.line.color.rgb = BORDER_LIGHT

    # 안내 텍스트
    guide_box = slide.shapes.add_textbox(
        Inches(4), Inches(3.8),
        Inches(5), Inches(1)
    )
    tf = guide_box.text_frame
    p = tf.paragraphs[0]
    p.text = "차트 이미지가 들어갈 영역"
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_MUTED
    p.alignment = PP_ALIGN.CENTER


def _add_table_slide(prs: Presentation) -> None:
    """
    표 슬라이드 추가 (6번 슬라이드)

    데이터 테이블을 표시하는 슬라이드입니다.
    샘플 데이터가 포함된 표가 생성됩니다.

    레이아웃:
    --------
        ┌─────────────────────────────────────┐
        │ 표 제목                             │  ← 파란색 헤더 바
        ├─────────────────────────────────────┤
        │                                     │
        │   ┌────┬────────┬────────┬─────┐   │
        │   │지역│사업체수│종사자수│증감률│   │  ← 표 헤더 (파란색)
        │   ├────┼────────┼────────┼─────┤   │
        │   │서울│1,234,567│5,678,901│+2.3%│  │  ← 데이터 행
        │   │부산│ 456,789│1,234,567│+1.5%│  │     (짝수 행 회색)
        │   │경북│ 234,567│ 890,123│+3.1%│  │
        │   │경남│ 345,678│1,012,345│+2.8%│  │
        │   │전체│5,678,901│15,234,567│+2.5%│ │
        │   └────┴────────┴────────┴─────┘   │
        │                                     │
        └─────────────────────────────────────┘

    Args:
        prs (Presentation): 프레젠테이션 객체

    스타일:
        - 표 헤더: #1243A6 배경, 흰색 텍스트, 14pt 볼드
        - 데이터 행: 12pt, 중앙 정렬
        - 짝수 행: #F5F7FA 배경
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 헤더
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_DARK
    header.line.fill.background()

    title = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.35),
        Inches(10), Inches(0.6)
    )
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "표 제목"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE

    # -------------------------------------------------------------------------
    # 샘플 표 생성
    # -------------------------------------------------------------------------
    rows, cols = 6, 5  # 헤더 + 5개 데이터 행
    table = slide.shapes.add_table(
        rows, cols,
        Inches(0.5), Inches(1.5),
        Inches(12.333), Inches(5)
    ).table

    # 열 너비 균등 분배
    col_width = Inches(12.333 / cols)
    for i in range(cols):
        table.columns[i].width = int(col_width)

    # 헤더 행
    headers = ["지역", "사업체수", "종사자수", "증감률", "비고"]
    for j, header_text in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header_text
        cell.fill.solid()
        cell.fill.fore_color.rgb = BG_TABLE_HEADER

        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        p.alignment = PP_ALIGN.CENTER

    # 데이터 행 (샘플)
    sample_data = [
        ["서울", "1,234,567", "5,678,901", "+2.3%", ""],
        ["부산", "456,789", "1,234,567", "+1.5%", ""],
        ["경북", "234,567", "890,123", "+3.1%", ""],
        ["경남", "345,678", "1,012,345", "+2.8%", ""],
        ["전체", "5,678,901", "15,234,567", "+2.5%", ""],
    ]

    for i, row_data in enumerate(sample_data):
        for j, value in enumerate(row_data):
            cell = table.cell(i + 1, j)
            cell.text = value

            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.alignment = PP_ALIGN.CENTER

            # 짝수 행 배경색 (가독성 향상)
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = BG_TABLE_ROW


def _add_insights_slide(prs: Presentation) -> None:
    """
    인사이트 슬라이드 추가 (7번 슬라이드)

    분석 결과로 도출된 인사이트를 카드 형태로 나열하는 슬라이드입니다.
    아이콘, 제목, 상세 내용으로 구성된 카드가 세로로 배치됩니다.

    레이아웃:
    --------
        ┌─────────────────────────────────────┐
        │ 주요 인사이트                        │  ← 파란색 헤더 바
        ├─────────────────────────────────────┤
        │                                     │
        │   ┌───────────────────────────────┐ │
        │   │ 📍 지역 특성                   │ │  ← 인사이트 카드 1
        │   │   경상북도는 제조업 비중이...   │ │
        │   └───────────────────────────────┘ │
        │                                     │
        │   ┌───────────────────────────────┐ │
        │   │ 📈 성장 추세                   │ │  ← 인사이트 카드 2
        │   │   전분기 대비 2.5% 증가...     │ │
        │   └───────────────────────────────┘ │
        │                                     │
        │   ┌───────────────────────────────┐ │
        │   │ ⚠️ 주의 사항                   │ │  ← 인사이트 카드 3
        │   │   일부 지역 폐업률 증가...     │ │
        │   └───────────────────────────────┘ │
        │                                     │
        └─────────────────────────────────────┘

    Args:
        prs (Presentation): 프레젠테이션 객체

    스타일:
        - 카드 배경: #F8F9FA (연한 회색)
        - 카드 테두리: #DCDCDC
        - 제목: 16pt, 볼드, 진한 회색
        - 내용: 14pt, 회색
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 헤더
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_DARK
    header.line.fill.background()

    title = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.35),
        Inches(10), Inches(0.6)
    )
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "주요 인사이트"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE

    # -------------------------------------------------------------------------
    # 인사이트 카드들
    # -------------------------------------------------------------------------
    insights = [
        ("📍 지역 특성", "경상북도는 제조업 비중이 높아 산업 다각화가 필요합니다."),
        ("📈 성장 추세", "전분기 대비 2.5% 증가하여 긍정적인 성장세를 보입니다."),
        ("⚠️ 주의 사항", "일부 지역의 폐업률이 증가 추세에 있어 모니터링이 필요합니다."),
    ]

    card_height = Inches(1.5)
    start_top = Inches(1.5)
    gap = Inches(0.2)

    for i, (title_text, content) in enumerate(insights):
        top = start_top + i * (card_height + gap)

        # 카드 배경 (둥근 사각형)
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), top,
            Inches(12.333), card_height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = BG_LIGHT
        card.line.color.rgb = BORDER_CARD

        # 제목 (아이콘 포함)
        title_box = slide.shapes.add_textbox(
            Inches(0.8), top + Inches(0.2),
            Inches(11.733), Inches(0.5)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(44, 62, 80)  # 진한 회색

        # 내용
        content_box = slide.shapes.add_textbox(
            Inches(1), top + Inches(0.7),
            Inches(11.533), Inches(0.7)
        )
        tf = content_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = content
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(85, 85, 85)  # 회색


def _add_thankyou_slide(prs: Presentation) -> None:
    """
    마무리 슬라이드 추가 (8번 슬라이드)

    보고서의 마지막 페이지로 사용되는 감사 인사 슬라이드입니다.
    제목 슬라이드와 유사한 디자인으로 일관성을 유지합니다.

    레이아웃:
    --------
        ┌─────────────────────────────────────┐
        │                                     │
        │                                     │
        │                                     │
        │            감사합니다                │  ← 48pt, 볼드, 흰색
        │                                     │
        │   문의: example@email.com           │  ← 연락처 정보
        │   담당: 홍길동                       │
        │                                     │
        │                                     │
        └─────────────────────────────────────┘

    Args:
        prs (Presentation): 프레젠테이션 객체

    스타일:
        - 배경: #1243A6 (진한 파랑)
        - 감사 문구: 48pt, 볼드, 흰색, 중앙 정렬
        - 연락처: 16pt, 연한 흰색
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 배경
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRIMARY_DARK
    bg.line.fill.background()

    # 감사 문구
    thanks_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.8),
        Inches(12.333), Inches(1.5)
    )
    tf = thanks_box.text_frame
    p = tf.paragraphs[0]
    p.text = "감사합니다"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.alignment = PP_ALIGN.CENTER

    # 연락처 정보
    contact_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(4.5),
        Inches(12.333), Inches(1)
    )
    tf = contact_box.text_frame
    p = tf.paragraphs[0]
    p.text = "문의: example@email.com | 담당: 홍길동"
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_LIGHT
    p.alignment = PP_ALIGN.CENTER


# =============================================================================
# 메인 실행 함수
# =============================================================================

def main() -> Path:
    """
    메인 실행 함수

    템플릿을 생성하고 templates 폴더에 저장합니다.
    실행 완료 후 결과를 콘솔에 출력합니다.

    Returns:
        Path: 생성된 템플릿 파일 경로

    사용 예:
    ------
        $ python module/create_ppt_template.py

    출력 예:
    ------
        ============================================================
        PPT 템플릿 생성 완료!
        ============================================================

        저장 위치: C:/project/templates/report_template.pptx

        포함된 슬라이드:
          1. 제목 슬라이드
          2. 목차 슬라이드
          ...

        이 템플릿을 PowerPoint에서 열어 마스터 슬라이드를 수정하세요.
        수정 방법: 보기 → 슬라이드 마스터
        ============================================================
    """
    # 템플릿 생성
    prs = create_template()

    # -------------------------------------------------------------------------
    # 저장 경로 설정
    # module 폴더의 부모(프로젝트 루트) → templates 폴더
    # -------------------------------------------------------------------------
    output_dir = Path(__file__).parent.parent / "templates"
    output_dir.mkdir(exist_ok=True)

    output_path = output_dir / "report_template.pptx"
    prs.save(str(output_path))

    # -------------------------------------------------------------------------
    # 결과 출력
    # -------------------------------------------------------------------------
    print("=" * 60)
    print("PPT 템플릿 생성 완료!")
    print("=" * 60)
    print(f"\n저장 위치: {output_path}")
    print(f"\n포함된 슬라이드:")
    print("  1. 제목 슬라이드")
    print("  2. 목차 슬라이드")
    print("  3. 내용 슬라이드 (제목+본문)")
    print("  4. 지표 카드 슬라이드 (4개 카드)")
    print("  5. 차트 슬라이드")
    print("  6. 표 슬라이드")
    print("  7. 인사이트 슬라이드")
    print("  8. 마무리 슬라이드")
    print("\n이 템플릿을 PowerPoint에서 열어 마스터 슬라이드를 수정하세요.")
    print("수정 방법: 보기 → 슬라이드 마스터")
    print("=" * 60)

    return output_path


# =============================================================================
# 스크립트 직접 실행 시
# =============================================================================

if __name__ == "__main__":
    main()
