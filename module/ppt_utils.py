# -*- coding: utf-8 -*-
"""
PowerPoint(PPT) 문서 생성 유틸리티
====================================

python-pptx를 사용하여 PowerPoint 문서를 생성합니다.
클라우드 서버(Linux)에서도 사용 가능합니다.

주요 기능:
    1. 슬라이드 생성
    2. 텍스트/표/차트 삽입
    3. 이미지 삽입

필수 패키지:
    pip install python-pptx

사용 예:
    from module.ppt_utils import create_dashboard_ppt

    ppt_path = create_dashboard_ppt(data_dict, charts_dict)

Author: Claude AI Agent
Created: 2024-12-18
"""

import os
import tempfile
from pathlib import Path
from datetime import datetime
from io import BytesIO

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Cm
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.chart import XL_CHART_TYPE
    PPT_AVAILABLE = True
except ImportError as e:
    PPT_AVAILABLE = False
    print(f"[WARNING] python-pptx import 오류: {e}")

import pandas as pd


def add_logo_to_slide(slide):
    """
    슬라이드에 경북 로고 추가

    Args:
        slide: 슬라이드 객체
    """
    logo_path = Path(__file__).parent.parent / "image" / "gyeongbuk_logo.png"
    if logo_path.exists():
        slide.shapes.add_picture(
            str(logo_path),
            Inches(0.2),   # 왼쪽 상단
            Inches(0.1),
            width=Inches(1.3)
        )


def create_title_slide(prs, title, subtitle=""):
    """
    제목 슬라이드 생성

    Args:
        prs: Presentation 객체
        title: 제목
        subtitle: 부제목

    Returns:
        slide: 생성된 슬라이드
    """
    slide_layout = prs.slide_layouts[6]  # 빈 슬라이드
    slide = prs.slides.add_slide(slide_layout)

    # 배경색 설정 (그라데이션 효과는 제한적이므로 단색 사용)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(102, 126, 234)  # #667eea

    # 제목 텍스트 박스
    left = Inches(0.5)
    top = Inches(2.5)
    width = Inches(9)
    height = Inches(1.5)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # 부제목
    if subtitle:
        top2 = Inches(4.2)
        txBox2 = slide.shapes.add_textbox(left, top2, width, Inches(0.8))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.color.rgb = RGBColor(255, 255, 255)
        p2.alignment = PP_ALIGN.CENTER

    # 로고 추가
    add_logo_to_slide(slide)

    return slide


def create_metrics_slide(prs, title, metrics):
    """
    주요 지표 슬라이드 생성 (카드 형식)

    Args:
        prs: Presentation 객체
        title: 슬라이드 제목
        metrics: 지표 리스트 [{'label': '라벨', 'value': '값', 'unit': '단위'}, ...]

    Returns:
        slide: 생성된 슬라이드
    """
    slide_layout = prs.slide_layouts[6]  # 빈 슬라이드
    slide = prs.slides.add_slide(slide_layout)

    # 제목 추가
    add_slide_title(slide, title)

    # 카드 배치 (2x2 그리드)
    card_width = Inches(4.2)
    card_height = Inches(2)
    start_left = Inches(0.5)
    start_top = Inches(1.5)
    gap = Inches(0.3)

    colors = [
        RGBColor(102, 126, 234),  # 파란색
        RGBColor(118, 75, 162),   # 보라색
        RGBColor(46, 204, 113),   # 초록색
        RGBColor(52, 152, 219),   # 하늘색
    ]

    for i, metric in enumerate(metrics[:4]):
        row = i // 2
        col = i % 2

        left = start_left + col * (card_width + gap)
        top = start_top + row * (card_height + gap)

        # 카드 배경
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, card_width, card_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors[i % len(colors)]
        shape.line.fill.background()

        # 라벨
        label_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.3), card_width - Inches(0.4), Inches(0.5))
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = metric.get('label', '')
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(255, 255, 255)

        # 값
        value_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.8), card_width - Inches(0.4), Inches(0.8))
        tf2 = value_box.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = str(metric.get('value', ''))
        p2.font.size = Pt(32)
        p2.font.bold = True
        p2.font.color.rgb = RGBColor(255, 255, 255)

        # 단위
        unit_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(1.5), card_width - Inches(0.4), Inches(0.4))
        tf3 = unit_box.text_frame
        p3 = tf3.paragraphs[0]
        p3.text = metric.get('unit', '')
        p3.font.size = Pt(12)
        p3.font.color.rgb = RGBColor(230, 230, 230)

    # 로고 추가
    add_logo_to_slide(slide)

    return slide


def create_table_slide(prs, title, df, max_rows=10):
    """
    표 슬라이드 생성

    Args:
        prs: Presentation 객체
        title: 슬라이드 제목
        df: pandas DataFrame
        max_rows: 최대 행 수

    Returns:
        slide: 생성된 슬라이드
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 제목 추가
    add_slide_title(slide, title)

    # 데이터 제한
    df_display = df.head(max_rows)
    rows = len(df_display) + 1  # 헤더 포함
    cols = len(df_display.columns)

    # 표 크기 계산
    left = Inches(0.3)
    top = Inches(1.4)
    width = Inches(9.4)
    height = Inches(0.4 * rows)

    # 표 생성
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # 열 너비 설정
    col_width = width / cols
    for i in range(cols):
        table.columns[i].width = int(col_width)

    # 헤더 설정
    for j, col_name in enumerate(df_display.columns):
        cell = table.cell(0, j)
        cell.text = str(col_name)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(102, 126, 234)

        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.size = Pt(10)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        paragraph.alignment = PP_ALIGN.CENTER

    # 데이터 설정
    for i, (idx, row) in enumerate(df_display.iterrows()):
        for j, value in enumerate(row):
            cell = table.cell(i + 1, j)

            # 값 포맷팅
            if pd.isna(value):
                text = '-'
            elif isinstance(value, (int, float)):
                if abs(value) >= 1000:
                    text = f'{value:,.0f}'
                elif isinstance(value, float):
                    text = f'{value:.2f}'
                else:
                    text = str(value)
            else:
                text = str(value)

            cell.text = text
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(9)
            paragraph.alignment = PP_ALIGN.CENTER

            # 짝수 행 배경색
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(245, 247, 250)

    # 로고 추가
    add_logo_to_slide(slide)

    return slide


def create_chart_slide(prs, title, image_path):
    """
    차트 이미지 슬라이드 생성

    Args:
        prs: Presentation 객체
        title: 슬라이드 제목
        image_path: 차트 이미지 경로

    Returns:
        slide: 생성된 슬라이드
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 제목 추가
    add_slide_title(slide, title)

    # 이미지 추가
    if os.path.exists(image_path):
        left = Inches(0.5)
        top = Inches(1.3)
        width = Inches(9)
        slide.shapes.add_picture(image_path, left, top, width=width)

    # 로고 추가
    add_logo_to_slide(slide)

    return slide


def create_insights_slide(prs, title, insights):
    """
    인사이트 슬라이드 생성

    Args:
        prs: Presentation 객체
        title: 슬라이드 제목
        insights: 인사이트 리스트 [{'icon': '아이콘', 'title': '제목', 'content': '내용'}, ...]

    Returns:
        slide: 생성된 슬라이드
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 제목 추가
    add_slide_title(slide, title)

    # 인사이트 카드 배치
    card_height = Inches(1.2)
    start_top = Inches(1.4)
    left = Inches(0.5)
    width = Inches(9)
    gap = Inches(0.15)

    for i, insight in enumerate(insights[:4]):
        top = start_top + i * (card_height + gap)

        # 카드 배경
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, width, card_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(248, 249, 250)
        shape.line.color.rgb = RGBColor(220, 220, 220)

        # 아이콘 + 제목
        title_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{insight.get('icon', '')} {insight.get('title', '')}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(44, 62, 80)

        # 내용
        content_box = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(0.55), width - Inches(0.6), Inches(0.6))
        tf2 = content_box.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = insight.get('content', '')
        p2.font.size = Pt(11)
        p2.font.color.rgb = RGBColor(85, 85, 85)

    # 로고 추가
    add_logo_to_slide(slide)

    return slide


def add_slide_title(slide, title):
    """
    슬라이드에 제목 추가

    Args:
        slide: 슬라이드 객체
        title: 제목 텍스트
    """
    left = Inches(0.3)
    top = Inches(0.3)
    width = Inches(9.4)
    height = Inches(0.8)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(44, 62, 80)


def create_dashboard_ppt(
    title,
    subtitle,
    metrics,
    df_aggregated,
    df_pop_biz,
    insights,
    chart_paths=None,
    output_path=None
):
    """
    대시보드 PPT 생성 메인 함수

    Args:
        title: 발표 제목
        subtitle: 부제목
        metrics: 주요 지표 리스트
        df_aggregated: 지역별 집계 데이터
        df_pop_biz: 인구 밀도 데이터
        insights: 인사이트 리스트
        chart_paths: 차트 이미지 경로 딕셔너리
        output_path: 저장 경로 (None이면 임시 파일)

    Returns:
        str: 생성된 PPT 파일 경로
    """
    if not PPT_AVAILABLE:
        raise RuntimeError("python-pptx가 설치되지 않았습니다. pip install python-pptx")

    # Presentation 생성
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 1. 제목 슬라이드
    create_title_slide(prs, title, subtitle)

    # 2. 주요 지표 슬라이드
    create_metrics_slide(prs, "주요 지표 요약", metrics)

    # 3. 지역별 현황 표
    if df_aggregated is not None and len(df_aggregated) > 0:
        df_table = df_aggregated[['지역명', '총사업체수', '총종사자수', 'HHI', '1인당매출액']].copy()
        create_table_slide(prs, "지역별 현황", df_table)

    # 4. 차트 슬라이드들
    if chart_paths:
        for chart_title, chart_path in chart_paths.items():
            if os.path.exists(chart_path):
                create_chart_slide(prs, chart_title, chart_path)

    # 5. 인구 밀도 표
    if df_pop_biz is not None and len(df_pop_biz) > 0:
        df_density = df_pop_biz[['시도명', '총인구', '사업체수', '인구천명당사업체수']].copy()
        create_table_slide(prs, "인구 대비 사업체 밀도", df_density)

    # 6. 인사이트 슬라이드
    if insights:
        create_insights_slide(prs, "주요 인사이트", insights)

    # 저장
    if output_path is None:
        temp_dir = tempfile.mkdtemp()
        output_path = os.path.join(temp_dir, "dashboard.pptx")

    prs.save(output_path)
    return output_path


# 테스트
if __name__ == "__main__":
    if PPT_AVAILABLE:
        print("python-pptx 사용 가능")

        # 테스트 데이터
        metrics = [
            {'label': '총 사업체수', 'value': '1,234,567', 'unit': '개'},
            {'label': '총 종사자수', 'value': '5,678,901', 'unit': '명'},
            {'label': '평균 HHI', 'value': '1,523.4', 'unit': '낮을수록 다양함'},
            {'label': '1인당 매출액', 'value': '125.6', 'unit': '백만원'},
        ]

        insights = [
            {'icon': '📍', 'title': '사업체 밀도 최고 지역', 'content': '서울이 인구 천명당 45.3개로 가장 높습니다.'},
            {'icon': '📊', 'title': '산업 다양성', 'content': '평균 HHI는 1523.4로 다양한 산업 구조를 보입니다.'},
        ]

        output = create_dashboard_ppt(
            title="기업통계등록부 분석 보고서",
            subtitle="2024년 4분기 | 경상북도",
            metrics=metrics,
            df_aggregated=None,
            df_pop_biz=None,
            insights=insights
        )
        print(f"PPT 생성 완료: {output}")
    else:
        print("python-pptx 사용 불가")
