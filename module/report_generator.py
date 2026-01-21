# -*- coding: utf-8 -*-
"""
대시보드 보고서 생성 모듈 (Dashboard Report Generator)
=====================================================

이 모듈은 대시보드 데이터를 마크다운(MD)과 PowerPoint(PPT) 형식으로 변환하는
범용 보고서 생성 기능을 제공합니다.

주요 기능:
---------
1. 대시보드 데이터 수집 및 구조화
2. 마크다운(MD) 형식 보고서 생성
3. PowerPoint(PPT) 형식 보고서 생성 (템플릿 지원)
4. PPT 템플릿 분석 유틸리티

사용 흐름 (Workflow):
-------------------
    대시보드 데이터 → DashboardReport 객체 → MD 파일 → PPT 파일

모듈 의존성:
----------
    - pathlib: 크로스 플랫폼 경로 처리
    - pandas: 데이터프레임 처리
    - python-pptx: PPT 생성 (선택적)

사용 예시:
--------
    >>> from module.report_generator import DashboardReport
    >>>
    >>> # 1. 보고서 객체 생성
    >>> report = DashboardReport(
    ...     title="기업통계등록부 분석 보고서",
    ...     subtitle="2024년 4분기 | 경상북도",
    ...     source_file=__file__  # 현재 파일 경로 (MD 저장 위치 결정)
    ... )
    >>>
    >>> # 2. 데이터 추가
    >>> report.add_metrics([
    ...     {'label': '총 사업체수', 'value': '1,234,567', 'unit': '개'},
    ...     {'label': '총 종사자수', 'value': '5,678,901', 'unit': '명'},
    ... ])
    >>> report.add_table('지역별 현황', df_aggregated)
    >>> report.add_chart('차트 제목', chart_image_path)
    >>> report.add_insights([
    ...     {'icon': '📍', 'title': '핵심 발견', 'content': '상세 내용...'}
    ... ])
    >>>
    >>> # 3. MD 저장 (소스 파일과 같은 폴더에 같은 이름으로)
    >>> md_path = report.save_markdown()
    >>> # 예: /project/routes/대시보드1.py → /project/routes/대시보드1.md
    >>>
    >>> # 4. PPT 저장 (템플릿 지정 가능)
    >>> ppt_path = report.save_ppt(template_path='templates/report_template.pptx')
    >>> # 예: /project/routes/대시보드1.pptx

Author: Claude AI Agent
Created: 2025-01-15
License: MIT
"""

# =============================================================================
# 표준 라이브러리 임포트
# =============================================================================
from pathlib import Path
from datetime import datetime
from dataclasses import dataclass, field
from typing import Optional, List, Dict, Any
import tempfile

# =============================================================================
# 서드파티 라이브러리 임포트
# =============================================================================
import pandas as pd

# =============================================================================
# PPT 관련 임포트 (선택적)
# python-pptx가 설치되지 않은 환경에서도 MD 생성은 가능
# =============================================================================
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    PPT_AVAILABLE = True
except ImportError:
    PPT_AVAILABLE = False


# =============================================================================
# 데이터 클래스 정의 (Data Classes)
# 보고서 구성 요소를 구조화하기 위한 데이터 컨테이너
# =============================================================================

@dataclass
class ChartData:
    """
    차트 데이터 구조 클래스

    차트 이미지와 관련 정보를 담는 데이터 컨테이너입니다.

    Attributes:
        title (str): 차트 제목 (예: "연령대별 인구 분포")
        image_path (str): 차트 이미지 파일 경로 (예: "/tmp/chart_12345.png")
        description (str): 차트에 대한 설명 (선택적)

    Example:
        >>> chart = ChartData(
        ...     title="고령화율 추이",
        ...     image_path="/charts/aging_rate.png",
        ...     description="2020-2024년 경상북도 고령화율 변화"
        ... )
    """
    title: str
    image_path: str
    description: str = ""


@dataclass
class TableData:
    """
    표 데이터 구조 클래스

    pandas DataFrame과 표시 옵션을 담는 데이터 컨테이너입니다.

    Attributes:
        title (str): 표 제목 (예: "시군구별 인구 현황")
        dataframe (pd.DataFrame): 표 데이터가 담긴 DataFrame
        max_rows (int): 표시할 최대 행 수 (기본값: 15)
                       PPT/MD에서 너무 긴 표를 방지

    Example:
        >>> table = TableData(
        ...     title="인구 상위 10개 시군구",
        ...     dataframe=df_population,
        ...     max_rows=10
        ... )
    """
    title: str
    dataframe: pd.DataFrame
    max_rows: int = 15


@dataclass
class MetricData:
    """
    지표 데이터 구조 클래스

    대시보드의 KPI 카드에 표시되는 주요 지표 정보를 담습니다.

    Attributes:
        label (str): 지표 라벨 (예: "총 사업체수")
        value (str): 지표 값 (예: "1,234,567")
        unit (str): 단위 (예: "개", "명", "%")
        description (str): 부가 설명 (선택적)

    Example:
        >>> metric = MetricData(
        ...     label="고령화율",
        ...     value="23.5",
        ...     unit="%",
        ...     description="65세 이상 인구 비율"
        ... )
    """
    label: str
    value: str
    unit: str = ""
    description: str = ""


@dataclass
class InsightData:
    """
    인사이트 데이터 구조 클래스

    분석 결과로 도출된 인사이트 정보를 담습니다.

    Attributes:
        icon (str): 아이콘 이모지 (예: "📍", "📈", "⚠️")
        title (str): 인사이트 제목 (예: "핵심 발견")
        content (str): 인사이트 상세 내용

    Example:
        >>> insight = InsightData(
        ...     icon="📍",
        ...     title="지역 특성",
        ...     content="경상북도는 전국 평균 대비 고령화율이 1.5배 높습니다."
        ... )
    """
    icon: str
    title: str
    content: str


# =============================================================================
# 메인 클래스: DashboardReport
# =============================================================================

class DashboardReport:
    """
    대시보드 보고서 생성 클래스

    대시보드의 다양한 데이터(지표, 표, 차트, 인사이트)를 수집하고
    마크다운(MD) 또는 PowerPoint(PPT) 형식으로 변환합니다.

    디자인 원칙:
    ----------
    1. 데이터 수집 단계와 출력 단계 분리
    2. 소스 파일 위치 기반 자동 출력 경로 결정
    3. PPT 템플릿 지원으로 일관된 스타일 유지

    주요 메서드:
    ----------
    - add_metric(), add_metrics(): 지표 추가
    - add_table(): 표 데이터 추가
    - add_chart(): 차트 이미지 추가
    - add_insight(), add_insights(): 인사이트 추가
    - to_markdown(): 마크다운 문자열 생성
    - save_markdown(): 마크다운 파일 저장
    - save_ppt(): PowerPoint 파일 저장

    Attributes:
        title (str): 보고서 제목
        subtitle (str): 부제목
        created_at (datetime): 생성 시간
        source_path (Path): 소스 파일 경로
        output_dir (Path): 출력 디렉토리
        base_name (str): 기본 파일명 (확장자 제외)
        metrics (List[MetricData]): 지표 목록
        tables (List[TableData]): 표 목록
        charts (List[ChartData]): 차트 목록
        insights (List[InsightData]): 인사이트 목록

    Example:
        >>> # 기본 사용
        >>> report = DashboardReport(
        ...     title="인구통계 분석 보고서",
        ...     subtitle="2024년 12월 | 경상북도",
        ...     source_file=__file__
        ... )
        >>>
        >>> # 데이터 추가
        >>> report.add_metrics([...])
        >>> report.add_table("현황표", df)
        >>>
        >>> # 저장
        >>> report.save_markdown()  # .md 파일 생성
        >>> report.save_ppt()       # .pptx 파일 생성
    """

    def __init__(
        self,
        title: str,
        subtitle: str = "",
        source_file: Optional[str] = None,
        output_dir: Optional[str] = None
    ):
        """
        DashboardReport 초기화

        Args:
            title (str): 보고서 제목
                        예: "기업통계등록부 분석 보고서"

            subtitle (str, optional): 부제목. 기준일, 지역 등을 표시
                                     예: "2024년 4분기 | 경상북도"

            source_file (str, optional): 소스 파일의 절대 경로
                                        일반적으로 __file__을 전달
                                        이 경로를 기준으로 MD/PPT 저장 위치 결정
                                        예: 대시보드1.py → 대시보드1.md, 대시보드1.pptx

            output_dir (str, optional): 출력 디렉토리 경로
                                       지정하지 않으면 source_file과 같은 폴더 사용
                                       source_file도 없으면 시스템 임시 폴더 사용

        Example:
            >>> # source_file 지정 (권장)
            >>> report = DashboardReport(
            ...     title="보고서",
            ...     source_file=__file__
            ... )
            >>> # → source_file 폴더에 같은 이름으로 저장
            >>>
            >>> # output_dir 직접 지정
            >>> report = DashboardReport(
            ...     title="보고서",
            ...     output_dir="/reports/output"
            ... )
        """
        # 기본 정보 저장
        self.title = title
        self.subtitle = subtitle
        self.created_at = datetime.now()

        # ---------------------------------------------------------------------
        # 소스 파일 경로 처리
        # source_file이 지정되면 해당 파일과 같은 폴더에 같은 이름으로 저장
        # ---------------------------------------------------------------------
        if source_file:
            self.source_path = Path(source_file)
            # output_dir이 지정되면 그것을 사용, 아니면 source_file의 폴더 사용
            self.output_dir = Path(output_dir) if output_dir else self.source_path.parent
            # 확장자를 제외한 파일명 (예: 대시보드1.py → 대시보드1)
            self.base_name = self.source_path.stem
        else:
            self.source_path = None
            # source_file이 없으면 임시 폴더 또는 지정된 output_dir 사용
            self.output_dir = Path(output_dir) if output_dir else Path(tempfile.gettempdir())
            self.base_name = "report"

        # ---------------------------------------------------------------------
        # 데이터 저장소 초기화
        # 각 타입별로 리스트로 관리
        # ---------------------------------------------------------------------
        self.metrics: List[MetricData] = []      # 주요 지표 (KPI 카드)
        self.tables: List[TableData] = []        # 표 데이터
        self.charts: List[ChartData] = []        # 차트 이미지
        self.insights: List[InsightData] = []    # 분석 인사이트
        self.custom_sections: List[Dict] = []    # 커스텀 섹션

    # =========================================================================
    # 데이터 추가 메서드 (Data Addition Methods)
    # =========================================================================

    def add_metric(self, label: str, value: str, unit: str = "", description: str = ""):
        """
        단일 지표 추가

        대시보드의 KPI 카드에 표시되는 주요 지표를 추가합니다.

        Args:
            label (str): 지표 라벨 (예: "총 사업체수")
            value (str): 지표 값 (예: "1,234,567")
            unit (str, optional): 단위 (예: "개", "명", "%")
            description (str, optional): 부가 설명

        Example:
            >>> report.add_metric(
            ...     label="고령화율",
            ...     value="23.5",
            ...     unit="%",
            ...     description="65세 이상 인구 비율"
            ... )
        """
        self.metrics.append(MetricData(label, value, unit, description))

    def add_metrics(self, metrics: List[Dict[str, Any]]):
        """
        여러 지표 일괄 추가

        딕셔너리 리스트 형태로 여러 지표를 한 번에 추가합니다.
        대시보드에서 수집한 지표 데이터를 그대로 전달할 때 유용합니다.

        Args:
            metrics (List[Dict]): 지표 딕셔너리 리스트
                                각 딕셔너리는 다음 키를 가질 수 있음:
                                - 'label' (필수): 지표 라벨
                                - 'value' (필수): 지표 값
                                - 'unit' (선택): 단위
                                - 'description' (선택): 설명

        Example:
            >>> report.add_metrics([
            ...     {'label': '총 사업체수', 'value': '1,234,567', 'unit': '개'},
            ...     {'label': '총 종사자수', 'value': '5,678,901', 'unit': '명'},
            ...     {'label': '평균 HHI', 'value': '1,523.4', 'unit': ''},
            ...     {'label': '증감률', 'value': '+2.5%', 'unit': '전분기 대비'},
            ... ])
        """
        for m in metrics:
            self.metrics.append(MetricData(
                label=m.get('label', ''),
                value=str(m.get('value', '')),
                unit=m.get('unit', ''),
                description=m.get('description', '')
            ))

    def add_table(self, title: str, df: pd.DataFrame, max_rows: int = 15):
        """
        표 데이터 추가

        pandas DataFrame을 보고서에 표로 추가합니다.
        DataFrame은 복사되어 저장되므로 원본 데이터에 영향 없음.

        Args:
            title (str): 표 제목 (예: "시군구별 인구 현황")
            df (pd.DataFrame): 표 데이터
            max_rows (int, optional): 표시할 최대 행 수 (기본값: 15)
                                     PPT/MD에서 너무 긴 표 방지

        Example:
            >>> # 기본 사용
            >>> report.add_table("상위 10개 지역", df_top10)
            >>>
            >>> # 행 수 제한
            >>> report.add_table("전체 현황", df_all, max_rows=20)

        Note:
            - DataFrame은 .copy()로 복사되어 저장
            - 원본 DataFrame을 나중에 수정해도 보고서에 영향 없음
        """
        self.tables.append(TableData(title, df.copy(), max_rows))

    def add_chart(self, title: str, image_path: str, description: str = ""):
        """
        차트 이미지 추가

        차트 이미지 파일 경로를 보고서에 추가합니다.
        matplotlib 등으로 생성한 차트를 저장 후 경로를 전달합니다.

        Args:
            title (str): 차트 제목 (예: "연령대별 인구 분포")
            image_path (str): 차트 이미지 파일 경로
                             절대 경로 권장
            description (str, optional): 차트에 대한 설명

        Example:
            >>> # matplotlib으로 차트 생성 후 저장
            >>> import matplotlib.pyplot as plt
            >>> fig, ax = plt.subplots()
            >>> ax.bar(x, y)
            >>> fig.savefig('/tmp/chart.png')
            >>> plt.close(fig)
            >>>
            >>> # 보고서에 추가
            >>> report.add_chart(
            ...     title="연령대별 인구",
            ...     image_path='/tmp/chart.png',
            ...     description="2024년 12월 기준"
            ... )

        Note:
            PPT 저장 시 이미지 파일이 존재해야 함
        """
        self.charts.append(ChartData(title, image_path, description))

    def add_insight(self, icon: str, title: str, content: str):
        """
        단일 인사이트 추가

        분석 결과로 도출된 인사이트를 추가합니다.

        Args:
            icon (str): 아이콘 이모지 (예: "📍", "📈", "⚠️", "💡")
            title (str): 인사이트 제목 (예: "핵심 발견")
            content (str): 인사이트 상세 내용

        Example:
            >>> report.add_insight(
            ...     icon="📍",
            ...     title="지역 특성",
            ...     content="경상북도는 전국 평균 대비 고령화율이 1.5배 높습니다."
            ... )

        Tip:
            추천 아이콘:
            - 📍 지역/위치 관련
            - 📈 성장/증가 관련
            - 📉 감소/하락 관련
            - ⚠️ 주의/경고 관련
            - 💡 아이디어/제안
            - ✅ 긍정적 결과
        """
        self.insights.append(InsightData(icon, title, content))

    def add_insights(self, insights: List[Dict[str, str]]):
        """
        여러 인사이트 일괄 추가

        딕셔너리 리스트 형태로 여러 인사이트를 한 번에 추가합니다.

        Args:
            insights (List[Dict]): 인사이트 딕셔너리 리스트
                                 각 딕셔너리는 다음 키를 가짐:
                                 - 'icon': 아이콘 이모지
                                 - 'title': 제목
                                 - 'content': 내용

        Example:
            >>> report.add_insights([
            ...     {'icon': '📍', 'title': '사업체 밀도', 'content': '서울이 가장 높습니다.'},
            ...     {'icon': '📊', 'title': '산업 다양성', 'content': 'HHI 지수가 낮아 다양합니다.'},
            ...     {'icon': '⚠️', 'title': '주의 사항', 'content': '일부 지역 폐업률 증가 중.'},
            ... ])
        """
        for i in insights:
            self.insights.append(InsightData(
                icon=i.get('icon', ''),
                title=i.get('title', ''),
                content=i.get('content', '')
            ))

    def add_section(self, title: str, content: str):
        """
        커스텀 섹션 추가

        표준 형식에 맞지 않는 커스텀 내용을 추가합니다.

        Args:
            title (str): 섹션 제목
            content (str): 섹션 내용 (마크다운 형식 가능)

        Example:
            >>> report.add_section(
            ...     title="분석 방법론",
            ...     content="본 분석은 다음의 방법론을 사용했습니다:\\n\\n1. ..."
            ... )
        """
        self.custom_sections.append({'title': title, 'content': content})

    # =========================================================================
    # 마크다운 생성 메서드 (Markdown Generation)
    # =========================================================================

    def to_markdown(self) -> str:
        """
        수집된 데이터를 마크다운 문자열로 변환

        보고서에 추가된 모든 데이터(지표, 표, 차트, 인사이트)를
        마크다운 형식의 문자열로 변환합니다.

        Returns:
            str: 마크다운 형식의 보고서 문자열

        생성되는 구조:
        -----------
            # 제목
            **부제목**
            _작성일: YYYY년 MM월 DD일 HH:MM_
            ---
            ## 주요 지표
            | 지표 | 값 | 단위 |
            ...
            ## 표 제목
            | 컬럼1 | 컬럼2 | ...
            ...
            ## 차트
            ![차트제목](이미지경로)
            ...
            ## 주요 인사이트
            ### 📍 인사이트 제목
            내용...

        Example:
            >>> md_content = report.to_markdown()
            >>> print(md_content[:200])
            # 기업통계등록부 분석 보고서

            **2024년 4분기 | 경상북도**

            _작성일: 2024년 12월 15일 14:30_
            ...
        """
        lines = []

        # -----------------------------------------------------------------
        # 헤더 영역: 제목, 부제목, 작성일
        # -----------------------------------------------------------------
        lines.append(f"# {self.title}")
        lines.append("")

        if self.subtitle:
            lines.append(f"**{self.subtitle}**")
            lines.append("")

        lines.append(f"_작성일: {self.created_at.strftime('%Y년 %m월 %d일 %H:%M')}_")
        lines.append("")
        lines.append("---")
        lines.append("")

        # -----------------------------------------------------------------
        # 커스텀 섹션 (AI 인사이트 등) - 헤더 바로 다음에 표시
        # 사용자 요청: 인사이트를 상단에 먼저 표시
        # -----------------------------------------------------------------
        for section in self.custom_sections:
            lines.append(f"## {section['title']}")
            lines.append("")
            lines.append(section['content'])
            lines.append("")

        # -----------------------------------------------------------------
        # 주요 지표 섹션
        # 마크다운 테이블 형식으로 출력
        # -----------------------------------------------------------------
        if self.metrics:
            lines.append("## 주요 지표")
            lines.append("")
            lines.append("| 지표 | 값 | 단위 |")
            lines.append("|------|-----|------|")
            for m in self.metrics:
                lines.append(f"| {m.label} | {m.value} | {m.unit} |")
            lines.append("")

        # -----------------------------------------------------------------
        # 표 데이터 섹션
        # 각 표를 마크다운 테이블로 변환
        # -----------------------------------------------------------------
        for table in self.tables:
            lines.append(f"## {table.title}")
            lines.append("")
            lines.append(self._df_to_markdown(table.dataframe, table.max_rows))
            lines.append("")

        # -----------------------------------------------------------------
        # 차트 섹션
        # 마크다운 이미지 문법 사용
        # -----------------------------------------------------------------
        if self.charts:
            lines.append("## 차트")
            lines.append("")
            for chart in self.charts:
                lines.append(f"### {chart.title}")
                lines.append("")
                lines.append(f"![{chart.title}]({chart.image_path})")
                if chart.description:
                    lines.append(f"\n_{chart.description}_")
                lines.append("")

        # -----------------------------------------------------------------
        # 인사이트 섹션
        # 아이콘과 함께 표시
        # -----------------------------------------------------------------
        if self.insights:
            lines.append("## 주요 인사이트")
            lines.append("")
            for insight in self.insights:
                lines.append(f"### {insight.icon} {insight.title}")
                lines.append("")
                lines.append(insight.content)
                lines.append("")

        # 커스텀 섹션은 이미 헤더 다음에 렌더링됨 (상단 표시)

        return "\n".join(lines)

    def _df_to_markdown(self, df: pd.DataFrame, max_rows: int = 15) -> str:
        """
        DataFrame을 마크다운 표로 변환

        pandas DataFrame을 마크다운 테이블 문법으로 변환합니다.
        숫자는 천 단위 구분자가 추가되고, 결측값은 '-'로 표시됩니다.

        Args:
            df (pd.DataFrame): 변환할 DataFrame
            max_rows (int): 최대 표시 행 수 (기본값: 15)

        Returns:
            str: 마크다운 테이블 문자열

        변환 규칙:
        --------
        - 정수/실수 1000 이상: 천 단위 구분자 추가 (예: 1,234,567)
        - 실수 1000 미만: 소수점 2자리까지 (예: 23.45)
        - 결측값(NaN): '-'로 표시
        - 기타: 문자열로 변환

        Example:
            생성되는 형식:
            | 지역 | 인구 | 비율 |
            |------|-----|------|
            | 서울 | 1,234,567 | 23.45 |
            | 부산 | 567,890 | 12.34 |
        """
        # max_rows만큼만 표시
        df_display = df.head(max_rows)

        # 헤더 행 생성
        headers = "| " + " | ".join(str(col) for col in df_display.columns) + " |"
        # 구분선 생성
        separator = "|" + "|".join(["------" for _ in df_display.columns]) + "|"

        # 데이터 행 생성
        rows = []
        for _, row in df_display.iterrows():
            formatted_values = []
            for val in row:
                # 값 포맷팅
                if pd.isna(val):
                    # 결측값 처리
                    formatted_values.append("-")
                elif isinstance(val, (int, float)):
                    # 숫자 포맷팅
                    if abs(val) >= 1000:
                        # 1000 이상: 천 단위 구분자
                        formatted_values.append(f"{val:,.0f}")
                    elif isinstance(val, float):
                        # 소수점 2자리
                        formatted_values.append(f"{val:.2f}")
                    else:
                        formatted_values.append(str(val))
                else:
                    # 문자열 등 기타
                    formatted_values.append(str(val))
            rows.append("| " + " | ".join(formatted_values) + " |")

        return "\n".join([headers, separator] + rows)

    def save_markdown(self, output_path: Optional[str] = None) -> Path:
        """
        마크다운 파일 저장

        생성된 마크다운 내용을 파일로 저장합니다.

        Args:
            output_path (str, optional): 저장 경로
                                        None이면 소스 파일과 같은 위치에 같은 이름으로 저장
                                        예: 대시보드1.py → 대시보드1.md

        Returns:
            Path: 저장된 파일의 경로

        Example:
            >>> # 자동 경로 (권장)
            >>> md_path = report.save_markdown()
            >>> print(md_path)  # /project/routes/대시보드1.md
            >>>
            >>> # 수동 경로 지정
            >>> md_path = report.save_markdown('/reports/custom_report.md')

        Note:
            - 부모 디렉토리가 없으면 자동 생성
            - UTF-8 인코딩으로 저장
        """
        # 출력 경로 결정
        if output_path:
            md_path = Path(output_path)
        else:
            # source_file과 같은 폴더에 같은 이름으로
            md_path = self.output_dir / f"{self.base_name}.md"

        # 디렉토리 생성 (없으면)
        md_path.parent.mkdir(parents=True, exist_ok=True)

        # 파일 저장 (UTF-8 인코딩)
        md_path.write_text(self.to_markdown(), encoding='utf-8')

        return md_path

    # =========================================================================
    # PPT 생성 메서드 (PowerPoint Generation)
    # =========================================================================

    def save_ppt(
        self,
        output_path: Optional[str] = None,
        template_path: Optional[str] = None,
        orientation: str = 'landscape'
    ) -> Path:
        """
        PowerPoint 파일 저장

        수집된 데이터를 PowerPoint 프레젠테이션으로 저장합니다.
        템플릿을 지정하면 해당 템플릿의 스타일을 적용합니다.

        Args:
            output_path (str, optional): 저장 경로
                                        None이면 소스 파일과 같은 위치에 같은 이름으로 저장
            template_path (str, optional): PPT 템플릿 파일 경로
                                          None이면 기본 스타일 사용
            orientation (str): 슬라이드 방향
                              'landscape' (횡/가로) - 기본값, 10" x 7.5"
                              'portrait' (종/세로) - 7.5" x 10"

        Returns:
            Path: 저장된 파일의 경로

        Raises:
            RuntimeError: python-pptx가 설치되지 않은 경우

        생성되는 슬라이드:
        ---------------
        1. 제목 슬라이드
        2. 주요 지표 슬라이드 (지표가 있는 경우)
        3. 표 슬라이드 (표당 1개)
        4. 차트 슬라이드 (차트당 1개)
        5. 인사이트 슬라이드 (인사이트가 있는 경우)

        Example:
            >>> # 기본 스타일 (가로)
            >>> ppt_path = report.save_ppt()
            >>>
            >>> # 세로 방향
            >>> ppt_path = report.save_ppt(orientation='portrait')
            >>>
            >>> # 템플릿 적용
            >>> ppt_path = report.save_ppt(
            ...     template_path='templates/report_template.pptx'
            ... )

        Note:
            - python-pptx 패키지 필요 (pip install python-pptx)
            - 템플릿 파일이 없으면 기본 스타일로 생성
        """
        # python-pptx 설치 확인
        if not PPT_AVAILABLE:
            raise RuntimeError(
                "python-pptx가 설치되지 않았습니다.\n"
                "설치: pip install python-pptx"
            )

        # 출력 경로 결정
        if output_path:
            ppt_path = Path(output_path)
        else:
            ppt_path = self.output_dir / f"{self.base_name}.pptx"

        # -----------------------------------------------------------------
        # 프레젠테이션 생성
        # 템플릿이 있으면 템플릿 기반, 없으면 새로 생성
        # -----------------------------------------------------------------
        if template_path and Path(template_path).exists():
            prs = Presentation(template_path)
        else:
            prs = Presentation()
            # 슬라이드 크기 설정 (방향에 따라)
            if orientation == 'portrait':
                # 세로 (종)
                prs.slide_width = Inches(7.5)
                prs.slide_height = Inches(10)
            else:
                # 가로 (횡) - 기본값
                prs.slide_width = Inches(10)
                prs.slide_height = Inches(7.5)

        # 로고 경로 설정 (각 슬라이드에 추가)
        self._logo_path = Path(__file__).parent.parent / "image" / "gyeongbuk_logo.png"
        # 현재 프레젠테이션 객체 저장 (로고 위치 계산용)
        self._current_prs = prs

        # -----------------------------------------------------------------
        # 슬라이드 생성
        # -----------------------------------------------------------------
        # 1. 제목 슬라이드
        self._add_title_slide(prs)

        # 2. 주요 지표 슬라이드
        if self.metrics:
            self._add_metrics_slide(prs)

        # 3. 표 슬라이드 (각 표당 1개)
        for table in self.tables:
            self._add_table_slide(prs, table)

        # 4. 차트 슬라이드 (각 차트당 1개)
        for chart in self.charts:
            self._add_chart_slide(prs, chart)

        # 5. 인사이트 슬라이드
        if self.insights:
            self._add_insights_slide(prs)

        # 6. 커스텀 섹션 슬라이드 (## 섹션별로 슬라이드 생성)
        if self.custom_sections:
            self._add_text_slides_from_sections(prs, self.custom_sections)

        # -----------------------------------------------------------------
        # 파일 저장
        # -----------------------------------------------------------------
        ppt_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(ppt_path))

        return ppt_path

    def _add_logo_to_slide(self, slide, prs=None):
        """
        슬라이드에 경북 로고 추가

        각 슬라이드의 오른쪽 상단에 로고 이미지를 추가합니다.

        Args:
            slide: 슬라이드 객체
            prs: 프레젠테이션 객체 (슬라이드 너비 계산용)
        """
        if hasattr(self, '_logo_path') and self._logo_path.exists():
            logo_width = Inches(1.2)
            # 슬라이드 너비 가져오기 (prs가 있으면 사용, 없으면 기본값)
            if hasattr(self, '_current_prs') and self._current_prs:
                slide_width = self._current_prs.slide_width.inches
            else:
                slide_width = 10  # 기본값
            # 오른쪽 상단에 배치
            left = Inches(slide_width) - logo_width - Inches(0.2)
            slide.shapes.add_picture(
                str(self._logo_path),
                left,
                Inches(0.15),
                width=logo_width
            )

    def _get_blank_layout(self, prs: Presentation):
        """
        빈 슬라이드 레이아웃 반환

        PPT 템플릿에서 빈 슬라이드 레이아웃을 가져옵니다.
        템플릿에 따라 빈 레이아웃의 인덱스가 다를 수 있습니다.

        Args:
            prs (Presentation): 프레젠테이션 객체

        Returns:
            SlideLayout: 빈 슬라이드 레이아웃

        Note:
            일반적으로 인덱스 6이 빈 슬라이드이나,
            템플릿에 따라 다를 수 있어 예외 처리 포함
        """
        try:
            return prs.slide_layouts[6]  # 일반적으로 빈 슬라이드
        except IndexError:
            return prs.slide_layouts[-1]  # 없으면 마지막 레이아웃

    def _add_slide_title(self, slide, title: str):
        """
        슬라이드에 제목 텍스트 박스 추가

        슬라이드 상단에 일관된 스타일의 제목을 추가합니다.

        Args:
            slide: 슬라이드 객체
            title (str): 제목 텍스트

        스타일:
            - 위치: 좌상단 (0.3", 0.3")
            - 크기: 9.4" x 0.8"
            - 폰트: 24pt, 볼드, 진한 회색(#2C3E50)
        """
        left = Inches(0.3)
        top = Inches(0.3)
        width = Inches(9.4)
        height = Inches(0.8)

        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RGBColor(44, 62, 80)  # 진한 회색

    def _add_title_slide(self, prs: Presentation):
        """
        제목 슬라이드 추가

        보고서의 첫 번째 슬라이드로, 제목과 부제목을 표시합니다.
        흰색 배경에 파란색 그라데이션 막대바 스타일입니다.

        Args:
            prs (Presentation): 프레젠테이션 객체
        """
        slide_layout = self._get_blank_layout(prs)
        slide = prs.slides.add_slide(slide_layout)

        # 슬라이드 크기 가져오기
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        # 배경색 설정 (흰색)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)

        # 파란색 그라데이션 막대바 (제목 배경)
        # 그라데이션 효과를 위해 여러 개의 사각형을 겹쳐서 표현
        bar_top = Inches(2.2)
        bar_height = Inches(1.6)
        bar_width_inch = slide_width.inches

        # 그라데이션 시뮬레이션: 여러 색상 단계로 막대 생성
        gradient_colors = [
            RGBColor(18, 67, 166),    # 진한 파랑 #1243A6
            RGBColor(40, 90, 180),
            RGBColor(60, 115, 200),
            RGBColor(80, 140, 215),
            RGBColor(100, 160, 230),
            RGBColor(130, 180, 240),  # 연한 파랑
        ]
        segment_width = bar_width_inch / len(gradient_colors)

        for i, color in enumerate(gradient_colors):
            segment = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(i * segment_width),
                bar_top,
                Inches(segment_width + 0.05),  # 약간 겹침
                bar_height
            )
            segment.fill.solid()
            segment.fill.fore_color.rgb = color
            segment.line.fill.background()  # 테두리 없음

        # 제목 텍스트 박스 (막대바 위에)
        txBox = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(2.5),
            slide_width - Inches(1),
            Inches(1.0)
        )
        tf = txBox.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = self.title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)  # 흰색
        p.alignment = PP_ALIGN.CENTER

        # 부제목 텍스트 박스 (막대바 아래)
        if self.subtitle:
            txBox2 = slide.shapes.add_textbox(
                Inches(0.5),
                Inches(4.2),
                slide_width - Inches(1),
                Inches(0.8)
            )
            tf2 = txBox2.text_frame
            p2 = tf2.paragraphs[0]
            p2.text = self.subtitle
            p2.font.size = Pt(18)
            p2.font.color.rgb = RGBColor(68, 68, 68)
            p2.alignment = PP_ALIGN.CENTER

        # 로고 추가 (오른쪽 상단)
        self._add_logo_to_slide(slide)

    def _add_metrics_slide(self, prs: Presentation):
        """
        주요 지표 슬라이드 추가

        KPI 카드 형태로 주요 지표를 표시하는 슬라이드를 추가합니다.
        최대 4개의 지표를 2x2 그리드로 배치합니다.

        Args:
            prs (Presentation): 프레젠테이션 객체

        레이아웃:
            - 2x2 그리드 (최대 4개 카드)
            - 각 카드: 4.2" x 2"
            - 색상: 파랑, 보라, 초록, 하늘색 순환
        """
        slide_layout = self._get_blank_layout(prs)
        slide = prs.slides.add_slide(slide_layout)

        self._add_slide_title(slide, "주요 지표 요약")

        # 카드 배치 설정
        card_width = Inches(4.2)
        card_height = Inches(2)
        start_left = Inches(0.5)
        start_top = Inches(1.5)
        gap = Inches(0.3)

        # 카드 색상 (순환)
        colors = [
            RGBColor(102, 126, 234),  # 파랑
            RGBColor(118, 75, 162),   # 보라
            RGBColor(46, 204, 113),   # 초록
            RGBColor(52, 152, 219),   # 하늘
        ]

        # 최대 4개 카드만 표시
        for i, metric in enumerate(self.metrics[:4]):
            row = i // 2  # 0 또는 1
            col = i % 2   # 0 또는 1

            left = start_left + col * (card_width + gap)
            top = start_top + row * (card_height + gap)

            # 카드 배경 (둥근 사각형)
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left, top, card_width, card_height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = colors[i % len(colors)]
            shape.line.fill.background()  # 테두리 없음

            # 라벨 텍스트
            label_box = slide.shapes.add_textbox(
                left + Inches(0.2), top + Inches(0.3),
                card_width - Inches(0.4), Inches(0.5)
            )
            tf = label_box.text_frame
            p = tf.paragraphs[0]
            p.text = metric.label
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(255, 255, 255)

            # 값 텍스트
            value_box = slide.shapes.add_textbox(
                left + Inches(0.2), top + Inches(0.8),
                card_width - Inches(0.4), Inches(0.8)
            )
            tf2 = value_box.text_frame
            p2 = tf2.paragraphs[0]
            p2.text = metric.value
            p2.font.size = Pt(32)
            p2.font.bold = True
            p2.font.color.rgb = RGBColor(255, 255, 255)

            # 단위 텍스트
            unit_box = slide.shapes.add_textbox(
                left + Inches(0.2), top + Inches(1.5),
                card_width - Inches(0.4), Inches(0.4)
            )
            tf3 = unit_box.text_frame
            p3 = tf3.paragraphs[0]
            p3.text = metric.unit
            p3.font.size = Pt(12)
            p3.font.color.rgb = RGBColor(230, 230, 230)

        # 로고 추가
        self._add_logo_to_slide(slide)

    def _add_table_slide(self, prs: Presentation, table_data: TableData):
        """
        표 슬라이드 추가

        DataFrame 데이터를 PPT 표로 변환한 슬라이드를 추가합니다.

        Args:
            prs (Presentation): 프레젠테이션 객체
            table_data (TableData): 표 데이터

        스타일:
            - 헤더 행: 파란색 배경, 흰색 텍스트, 볼드
            - 데이터 행: 짝수 행에 연한 회색 배경
            - 숫자: 천 단위 구분자, 소수점 2자리
            - 결측값: '-'로 표시
        """
        slide_layout = self._get_blank_layout(prs)
        slide = prs.slides.add_slide(slide_layout)

        self._add_slide_title(slide, table_data.title)

        df = table_data.dataframe.head(table_data.max_rows)
        rows = len(df) + 1  # 데이터 행 + 헤더 행
        cols = len(df.columns)

        # 표 위치 및 크기
        left = Inches(0.3)
        top = Inches(1.4)
        width = Inches(9.4)
        height = Inches(0.4 * rows)  # 행 수에 비례

        # 표 생성
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        # 열 너비 균등 분배
        col_width = width / cols
        for i in range(cols):
            table.columns[i].width = int(col_width)

        # 헤더 행 스타일
        for j, col_name in enumerate(df.columns):
            cell = table.cell(0, j)
            cell.text = str(col_name)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(102, 126, 234)  # 파란색
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE  # 세로 중간 정렬

            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(10)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(255, 255, 255)
            paragraph.alignment = PP_ALIGN.CENTER

        # 데이터 행
        for i, (idx, row) in enumerate(df.iterrows()):
            for j, value in enumerate(row):
                cell = table.cell(i + 1, j)

                # 값 포맷팅
                is_numeric = False
                if pd.isna(value):
                    text = '-'
                elif isinstance(value, (int, float)):
                    is_numeric = True
                    if abs(value) >= 1000:
                        text = f'{value:,.0f}'
                    elif isinstance(value, float):
                        text = f'{value:.2f}'
                    else:
                        text = str(value)
                else:
                    text = str(value)

                cell.text = text
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE  # 세로 중간 정렬
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.size = Pt(9)
                paragraph.alignment = PP_ALIGN.CENTER

                # 숫자는 볼드 처리
                if is_numeric:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = RGBColor(44, 62, 80)  # 진한 색

                # 짝수 행 배경색 (가독성 향상)
                if i % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(245, 247, 250)

        # 로고 추가
        self._add_logo_to_slide(slide)

    def _add_chart_slide(self, prs: Presentation, chart_data: ChartData):
        """
        차트 슬라이드 추가

        차트 이미지를 포함한 슬라이드를 추가합니다.

        Args:
            prs (Presentation): 프레젠테이션 객체
            chart_data (ChartData): 차트 데이터

        Note:
            이미지 파일이 존재해야 함
            이미지 너비는 9"로 고정, 높이는 비율 유지
        """
        slide_layout = self._get_blank_layout(prs)
        slide = prs.slides.add_slide(slide_layout)

        self._add_slide_title(slide, chart_data.title)

        # 이미지 추가
        image_path = Path(chart_data.image_path)
        if image_path.exists():
            left = Inches(0.5)
            top = Inches(1.3)
            width = Inches(9)
            slide.shapes.add_picture(str(image_path), left, top, width=width)

        # 로고 추가
        self._add_logo_to_slide(slide)

    def _add_insights_slide(self, prs: Presentation):
        """
        인사이트 슬라이드 추가

        분석 인사이트를 카드 형태로 나열한 슬라이드를 추가합니다.

        Args:
            prs (Presentation): 프레젠테이션 객체

        레이아웃:
            - 세로 나열 (최대 4개)
            - 각 카드: 높이 1.2", 너비 9"
            - 연한 회색 배경, 테두리
        """
        slide_layout = self._get_blank_layout(prs)
        slide = prs.slides.add_slide(slide_layout)

        self._add_slide_title(slide, "주요 인사이트")

        card_height = Inches(1.2)
        start_top = Inches(1.4)
        left = Inches(0.5)
        width = Inches(9)
        gap = Inches(0.15)

        # 최대 4개 인사이트만 표시
        for i, insight in enumerate(self.insights[:4]):
            top = start_top + i * (card_height + gap)

            # 카드 배경 (둥근 사각형)
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left, top, width, card_height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(248, 249, 250)  # 연한 회색
            shape.line.color.rgb = RGBColor(220, 220, 220)       # 테두리

            # 제목 (아이콘 + 제목)
            title_box = slide.shapes.add_textbox(
                left + Inches(0.2), top + Inches(0.15),
                width - Inches(0.4), Inches(0.4)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"{insight.icon} {insight.title}"
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = RGBColor(44, 62, 80)

            # 내용
            content_box = slide.shapes.add_textbox(
                left + Inches(0.3), top + Inches(0.55),
                width - Inches(0.6), Inches(0.6)
            )
            tf2 = content_box.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            p2.text = insight.content
            p2.font.size = Pt(11)
            p2.font.color.rgb = RGBColor(85, 85, 85)

        # 로고 추가
        self._add_logo_to_slide(slide)

    def _add_text_slide(self, prs: Presentation, title: str, content: str):
        """
        텍스트 내용 슬라이드 추가

        긴 텍스트 내용을 슬라이드에 표시합니다.
        내용이 길면 자동으로 잘라서 표시합니다.

        Args:
            prs (Presentation): 프레젠테이션 객체
            title (str): 슬라이드 제목
            content (str): 텍스트 내용
        """
        slide_layout = self._get_blank_layout(prs)
        slide = prs.slides.add_slide(slide_layout)

        self._add_slide_title(slide, title)

        # 텍스트 내용 영역
        left = Inches(0.5)
        top = Inches(1.4)
        width = Inches(9)
        height = Inches(5.5)

        # 마크다운 형식 정리
        import re
        clean_content = content
        clean_content = re.sub(r'^#{1,6}\s*', '', clean_content, flags=re.MULTILINE)  # # 제목 제거
        clean_content = re.sub(r'\*\*([^*]+)\*\*', r'\1', clean_content)  # **bold** -> bold
        clean_content = re.sub(r'\*([^*]+)\*', r'\1', clean_content)  # *italic* -> italic
        clean_content = re.sub(r'`([^`]+)`', r'\1', clean_content)  # `code` -> code
        clean_content = re.sub(r'```[\s\S]*?```', '[코드 블록 생략]', clean_content)  # 코드블록 제거

        # 내용이 너무 길면 자름 (약 1500자)
        if len(clean_content) > 1500:
            clean_content = clean_content[:1500] + "\n\n... (계속)"

        # 텍스트 박스 추가
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = clean_content
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(44, 62, 80)
        p.line_spacing = 1.3

        # 로고 추가
        self._add_logo_to_slide(slide)

    def _add_text_slides_from_sections(self, prs: Presentation, sections: list):
        """
        마크다운 섹션들을 여러 슬라이드로 분할하여 추가

        Args:
            prs (Presentation): 프레젠테이션 객체
            sections (list): [{'title': '제목', 'content': '내용'}, ...] 형태
        """
        for section in sections:
            title = section.get('title', '분석 결과')
            content = section.get('content', '')

            if not content.strip():
                continue

            # 내용이 매우 긴 경우 여러 슬라이드로 분할
            max_content_per_slide = 1200
            if len(content) > max_content_per_slide:
                # 단락 단위로 분할
                paragraphs = content.split('\n\n')
                current_content = ""
                slide_num = 1

                for para in paragraphs:
                    if len(current_content) + len(para) > max_content_per_slide:
                        # 현재까지의 내용으로 슬라이드 생성
                        if current_content.strip():
                            slide_title = f"{title} ({slide_num})" if slide_num > 1 else title
                            self._add_text_slide(prs, slide_title, current_content)
                            slide_num += 1
                        current_content = para
                    else:
                        current_content += "\n\n" + para if current_content else para

                # 남은 내용 처리
                if current_content.strip():
                    slide_title = f"{title} ({slide_num})" if slide_num > 1 else title
                    self._add_text_slide(prs, slide_title, current_content)
            else:
                self._add_text_slide(prs, title, content)


# =============================================================================
# 유틸리티 함수 (Utility Functions)
# PPT 템플릿 분석 및 디버깅용
# =============================================================================

def analyze_template(template_path: str) -> Dict[str, Any]:
    """
    PPT 템플릿의 레이아웃 구조 분석

    템플릿 파일의 슬라이드 크기와 레이아웃 정보를 분석합니다.
    새 템플릿을 적용하기 전에 구조를 파악하는 데 유용합니다.

    Args:
        template_path (str): 템플릿 파일 경로

    Returns:
        dict: 템플릿 정보
            - slide_width (float): 슬라이드 너비 (인치)
            - slide_height (float): 슬라이드 높이 (인치)
            - layouts (List[dict]): 레이아웃 목록
                - index (int): 레이아웃 인덱스
                - name (str): 레이아웃 이름
                - placeholders (List[dict]): 플레이스홀더 목록

    Raises:
        RuntimeError: python-pptx가 설치되지 않은 경우

    Example:
        >>> info = analyze_template('template.pptx')
        >>> print(f"크기: {info['slide_width']}\" x {info['slide_height']}\"")
        >>> for layout in info['layouts']:
        ...     print(f"[{layout['index']}] {layout['name']}")
    """
    if not PPT_AVAILABLE:
        raise RuntimeError("python-pptx가 설치되지 않았습니다.")

    prs = Presentation(template_path)

    info = {
        'slide_width': prs.slide_width.inches,
        'slide_height': prs.slide_height.inches,
        'layouts': []
    }

    for i, layout in enumerate(prs.slide_layouts):
        layout_info = {
            'index': i,
            'name': layout.name,
            'placeholders': []
        }

        for placeholder in layout.placeholders:
            layout_info['placeholders'].append({
                'idx': placeholder.placeholder_format.idx,
                'type': str(placeholder.placeholder_format.type),
                'name': placeholder.name
            })

        info['layouts'].append(layout_info)

    return info


def print_template_info(template_path: str):
    """
    템플릿 정보를 콘솔에 출력

    analyze_template()의 결과를 보기 좋게 포맷팅하여 출력합니다.
    템플릿 구조 파악 및 디버깅에 유용합니다.

    Args:
        template_path (str): 템플릿 파일 경로

    Example:
        >>> print_template_info('template.pptx')
        ============================================================
        템플릿 분석: template.pptx
        ============================================================
        슬라이드 크기: 13.33" x 7.50"

        레이아웃 목록:

          [0] Title Slide
              - idx=0, type=TITLE, name=Title 1
              - idx=1, type=SUBTITLE, name=Subtitle 2

          [1] Title and Content
              ...
    """
    info = analyze_template(template_path)

    print(f"\n{'='*60}")
    print(f"템플릿 분석: {template_path}")
    print(f"{'='*60}")
    print(f"슬라이드 크기: {info['slide_width']:.2f}\" x {info['slide_height']:.2f}\"")
    print(f"\n레이아웃 목록:")

    for layout in info['layouts']:
        print(f"\n  [{layout['index']}] {layout['name']}")
        for ph in layout['placeholders']:
            print(f"      - idx={ph['idx']}, type={ph['type']}, name={ph['name']}")


# =============================================================================
# 모듈 테스트 (Module Test)
# python -m module.report_generator 로 직접 실행 시 테스트
# =============================================================================

if __name__ == "__main__":
    print("=" * 60)
    print("DashboardReport 모듈 테스트")
    print("=" * 60)

    # 테스트 데이터로 보고서 생성
    report = DashboardReport(
        title="기업통계등록부 분석 보고서",
        subtitle="2024년 4분기 | 경상북도",
        source_file=__file__  # 이 파일 경로 기준으로 저장
    )

    # 지표 추가
    report.add_metrics([
        {'label': '총 사업체수', 'value': '1,234,567', 'unit': '개'},
        {'label': '총 종사자수', 'value': '5,678,901', 'unit': '명'},
        {'label': '평균 HHI', 'value': '1,523.4', 'unit': '낮을수록 다양함'},
        {'label': '1인당 매출액', 'value': '125.6', 'unit': '백만원'},
    ])

    # 인사이트 추가
    report.add_insights([
        {'icon': '📍', 'title': '사업체 밀도 최고', 'content': '서울이 인구 천명당 45.3개로 가장 높습니다.'},
        {'icon': '📊', 'title': '산업 다양성', 'content': '평균 HHI는 1523.4로 다양한 산업 구조입니다.'},
    ])

    # 마크다운 생성 테스트
    md_content = report.to_markdown()
    print("\n=== 생성된 마크다운 (처음 500자) ===")
    print(md_content[:500] + "...")

    # 마크다운 저장
    md_path = report.save_markdown()
    print(f"\n✅ MD 저장 완료: {md_path}")

    # PPT 저장
    if PPT_AVAILABLE:
        ppt_path = report.save_ppt()
        print(f"✅ PPT 저장 완료: {ppt_path}")
    else:
        print("⚠️ PPT 저장 불가 (python-pptx 미설치)")

    print("\n" + "=" * 60)
